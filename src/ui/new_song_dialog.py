"""Dialog for creating a new song from lyrics."""

import os
import re
from typing import Optional

from PyQt6.QtWidgets import (
    QDialog,
    QVBoxLayout,
    QHBoxLayout,
    QFormLayout,
    QLineEdit,
    QTextEdit,
    QPushButton,
    QLabel,
    QDialogButtonBox,
    QGroupBox,
    QMessageBox,
    QComboBox,
)
from PyQt6.QtCore import Qt

from ..models import Settings
from ..i18n import tr


class NewSongDialog(QDialog):
    """Dialog for creating a new song from title and lyrics."""

    def __init__(self, settings: Settings, base_path: str, parent=None):
        super().__init__(parent)
        self.settings = settings
        self.base_path = base_path
        self._created_folder: Optional[str] = None
        self._created_pptx: Optional[str] = None

        self._setup_ui()
        self._connect_signals()

    def _setup_ui(self) -> None:
        """Setup the dialog UI."""
        self.setWindowTitle(tr("dialog.newsong.title"))
        self.setMinimumSize(500, 500)
        self.resize(600, 600)

        layout = QVBoxLayout(self)

        # Song info group
        info_group = QGroupBox(tr("dialog.newsong.info"))
        info_layout = QFormLayout(info_group)

        self.title_input = QLineEdit()
        self.title_input.setPlaceholderText(tr("dialog.newsong.title_placeholder"))
        info_layout.addRow(tr("dialog.newsong.song_title"), self.title_input)

        # Subfolder selection
        self.subfolder_input = QComboBox()
        self.subfolder_input.setEditable(True)
        self._populate_subfolders()
        info_layout.addRow(tr("dialog.newsong.subfolder"), self.subfolder_input)

        layout.addWidget(info_group)

        # Lyrics group
        lyrics_group = QGroupBox(tr("dialog.newsong.lyrics"))
        lyrics_layout = QVBoxLayout(lyrics_group)

        self.lyrics_label = QLabel(tr("dialog.newsong.lyrics_hint"))
        self.lyrics_label.setWordWrap(True)
        lyrics_layout.addWidget(self.lyrics_label)

        self.lyrics_input = QTextEdit()
        self.lyrics_input.setPlaceholderText(tr("dialog.newsong.lyrics_placeholder"))
        lyrics_layout.addWidget(self.lyrics_input)

        layout.addWidget(lyrics_group)

        # Preview info
        self.preview_label = QLabel()
        layout.addWidget(self.preview_label)

        # Button box
        self.button_box = QDialogButtonBox(
            QDialogButtonBox.StandardButton.Ok | QDialogButtonBox.StandardButton.Cancel
        )
        self.button_box.button(QDialogButtonBox.StandardButton.Ok).setText(tr("dialog.newsong.create"))
        self.button_box.button(QDialogButtonBox.StandardButton.Cancel).setText(tr("button.cancel"))
        layout.addWidget(self.button_box)

        # Initially disable OK
        self.button_box.button(QDialogButtonBox.StandardButton.Ok).setEnabled(False)

    def _populate_subfolders(self) -> None:
        """Populate subfolder combo with existing song categories."""
        songs_path = self.settings.get_songs_path(self.base_path)
        self.subfolder_input.addItem("")  # Root of songs folder

        if os.path.isdir(songs_path):
            for item in sorted(os.listdir(songs_path)):
                item_path = os.path.join(songs_path, item)
                if os.path.isdir(item_path):
                    # Check if it's a category folder (contains subfolders, not just files)
                    has_subfolders = any(
                        os.path.isdir(os.path.join(item_path, sub))
                        for sub in os.listdir(item_path)
                    )
                    if has_subfolders:
                        self.subfolder_input.addItem(item)

    def _connect_signals(self) -> None:
        """Connect widget signals."""
        self.title_input.textChanged.connect(self._update_preview)
        self.lyrics_input.textChanged.connect(self._update_preview)
        self.button_box.accepted.connect(self._on_accept)
        self.button_box.rejected.connect(self.reject)

    def _update_preview(self) -> None:
        """Update the preview info."""
        title = self.title_input.text().strip()
        lyrics = self.lyrics_input.toPlainText().strip()

        if not title or not lyrics:
            self.preview_label.setText("")
            self.button_box.button(QDialogButtonBox.StandardButton.Ok).setEnabled(False)
            return

        # Count slides (separated by blank lines)
        slides = self._split_lyrics(lyrics)
        self.preview_label.setText(tr("dialog.newsong.preview", count=len(slides)))
        self.button_box.button(QDialogButtonBox.StandardButton.Ok).setEnabled(True)

    def _split_lyrics(self, lyrics: str) -> list:
        """Split lyrics into slides by blank lines."""
        # Normalize line endings
        lyrics = lyrics.replace('\r\n', '\n').replace('\r', '\n')

        # Split by one or more blank lines
        parts = re.split(r'\n\s*\n', lyrics)

        # Filter out empty parts and strip whitespace
        slides = [part.strip() for part in parts if part.strip()]

        return slides

    def _on_accept(self) -> None:
        """Handle dialog acceptance - create the song."""
        title = self.title_input.text().strip()
        lyrics = self.lyrics_input.toPlainText().strip()
        subfolder = self.subfolder_input.currentText().strip()

        if not title or not lyrics:
            return

        # Create safe folder name from title
        safe_name = self._make_safe_filename(title)

        # Build folder path
        songs_path = self.settings.get_songs_path(self.base_path)
        if subfolder:
            folder_path = os.path.join(songs_path, subfolder, safe_name)
        else:
            folder_path = os.path.join(songs_path, safe_name)

        # Check if folder already exists
        if os.path.exists(folder_path):
            result = QMessageBox.question(
                self,
                tr("dialog.newsong.title"),
                tr("dialog.newsong.folder_exists", name=safe_name),
                QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No
            )
            if result != QMessageBox.StandardButton.Yes:
                return

        try:
            # Create folder
            os.makedirs(folder_path, exist_ok=True)

            # Create song.properties
            props_path = os.path.join(folder_path, "song.properties")
            with open(props_path, "w", encoding="utf-8") as f:
                f.write(f"title={title}\n")

            # Create PowerPoint
            pptx_path = os.path.join(folder_path, f"{safe_name}.pptx")
            self._create_pptx(title, lyrics, pptx_path)

            self._created_folder = folder_path
            self._created_pptx = pptx_path

            self.accept()

        except Exception as e:
            QMessageBox.critical(
                self,
                tr("dialog.newsong.title"),
                tr("dialog.newsong.error", error=str(e))
            )

    def _make_safe_filename(self, name: str) -> str:
        """Convert a name to a safe filename."""
        # Remove or replace invalid characters
        safe = re.sub(r'[<>:"/\\|?*]', '', name)
        # Replace multiple spaces with single space
        safe = re.sub(r'\s+', ' ', safe)
        return safe.strip()

    def _create_pptx(self, title: str, lyrics: str, output_path: str) -> None:
        """Create a PowerPoint presentation from lyrics."""
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN

        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
        prs.slide_height = Inches(7.5)

        # Get blank layout
        blank_layout = prs.slide_layouts[6]  # Blank layout

        # Split lyrics into slides
        slides_text = self._split_lyrics(lyrics)

        for slide_text in slides_text:
            slide = prs.slides.add_slide(blank_layout)

            # Add title text box at top
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.3),
                Inches(12.333), Inches(0.8)
            )
            title_frame = title_box.text_frame
            title_para = title_frame.paragraphs[0]
            title_para.text = title
            title_para.font.size = Pt(28)
            title_para.font.bold = True
            title_para.alignment = PP_ALIGN.CENTER

            # Add lyrics text box
            lyrics_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.5),
                Inches(12.333), Inches(5.5)
            )
            lyrics_frame = lyrics_box.text_frame
            lyrics_frame.word_wrap = True

            # Add lyrics text
            lyrics_para = lyrics_frame.paragraphs[0]
            lyrics_para.text = slide_text
            lyrics_para.font.size = Pt(32)
            lyrics_para.alignment = PP_ALIGN.CENTER

        prs.save(output_path)

    def get_created_folder(self) -> Optional[str]:
        """Get the path to the created song folder."""
        return self._created_folder

    def get_created_pptx(self) -> Optional[str]:
        """Get the path to the created PowerPoint file."""
        return self._created_pptx
