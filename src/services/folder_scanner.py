"""Service for scanning song and generic folders."""

import os
import re
from typing import List, Optional, Dict

from pptx import Presentation

from ..models import Song, GenericItem, OfferingSlide, Settings
from ..logging_config import get_logger

logger = get_logger("folder_scanner")


class FolderScanner:
    """Scans folders for songs and generic items."""

    def __init__(self, settings: Settings, base_path: str = "."):
        self.settings = settings
        self.base_path = base_path
        self._songs_cache: Optional[List[Song]] = None
        self._generic_cache: Optional[List[GenericItem]] = None
        self._offerings_cache: Optional[List[OfferingSlide]] = None

    def refresh(self) -> None:
        """Force rescan by clearing all caches."""
        self._songs_cache = None
        self._generic_cache = None
        self._offerings_cache = None

    def scan_songs(self, force_refresh: bool = False) -> List[Song]:
        """
        Scan the Songs folder recursively.
        Returns a list of Song objects for all leaf folders containing song files.
        """
        if not force_refresh and self._songs_cache is not None:
            return self._songs_cache

        songs_path = self.settings.get_songs_path(self.base_path)
        if not os.path.exists(songs_path):
            self._songs_cache = []
            return []

        songs = []
        self._scan_songs_recursive(songs_path, songs_path, songs)
        self._songs_cache = sorted(songs, key=lambda s: s.relative_path.lower())
        return self._songs_cache

    def _scan_songs_recursive(
        self, current_path: str, songs_root: str, songs: List[Song]
    ) -> bool:
        """
        Recursively scan for songs.
        Returns True if this folder is a leaf (contains song files).
        """
        if not os.path.isdir(current_path):
            return False

        # Check if this folder contains any pptx or pdf files
        has_song_files = False
        has_subdirs = False

        for item in os.listdir(current_path):
            item_path = os.path.join(current_path, item)
            if os.path.isdir(item_path):
                has_subdirs = True
            elif item.lower().endswith((".pptx", ".ppt", ".pdf")):
                has_song_files = True

        # If this folder has song files, treat it as a leaf
        if has_song_files:
            song = Song.from_folder(current_path, songs_root)
            if song:
                songs.append(song)
            return True

        # Otherwise, recurse into subdirectories
        if has_subdirs:
            for item in os.listdir(current_path):
                item_path = os.path.join(current_path, item)
                if os.path.isdir(item_path):
                    self._scan_songs_recursive(item_path, songs_root, songs)

        return False

    def scan_generic(self, force_refresh: bool = False) -> List[GenericItem]:
        """
        Scan the Generic (Algemeen) folder (flat, ignoring hierarchy).
        Returns a list of GenericItem objects, excluding the Collecte file.
        """
        if not force_refresh and self._generic_cache is not None:
            return self._generic_cache

        generic_path = self.settings.get_algemeen_path(self.base_path)
        if not os.path.exists(generic_path):
            self._generic_cache = []
            return []

        items = []
        collecte_name = self.settings.collecte_filename.lower()
        stub_name = self.settings.stub_template_filename.lower()

        for item in os.listdir(generic_path):
            # Skip Collecte and StubTemplate
            if item.lower() in (collecte_name, stub_name):
                continue

            item_path = os.path.join(generic_path, item)

            # Only include PowerPoint files
            if os.path.isfile(item_path) and item.lower().endswith((".pptx", ".ppt")):
                items.append(GenericItem(name=item, pptx_path=item_path))

        self._generic_cache = sorted(items, key=lambda i: i.display_name.lower())
        return self._generic_cache

    # Backwards compatibility alias
    def scan_algemeen(self, force_refresh: bool = False) -> List[GenericItem]:
        """Scan the Generic folder. (Backwards compatibility alias)"""
        return self.scan_generic(force_refresh)

    def get_offering_slides(self, pptx_path: Optional[str] = None, force_refresh: bool = False) -> List[OfferingSlide]:
        """
        Get all slides from an Offering (Collecte) PowerPoint.
        Returns a list of OfferingSlide objects with index and title.

        Args:
            pptx_path: Path to the offering pptx file. If None, uses default Collecte path.
            force_refresh: If True, forces rescan even if cached.
        """
        # Use default collecte path if none provided
        if pptx_path is None:
            pptx_path = self.settings.get_collecte_path(self.base_path)
            # Only use cache for default collecte file
            if not force_refresh and self._offerings_cache is not None:
                return self._offerings_cache

        if not os.path.exists(pptx_path):
            logger.debug(f"Offering file not found: {pptx_path}")
            return []

        slides = []
        try:
            logger.debug(f"Loading offering slides from: {pptx_path}")
            prs = Presentation(pptx_path)
            for idx, slide in enumerate(prs.slides):
                title = self._extract_slide_title(slide, idx)
                slides.append(OfferingSlide(index=idx, title=title))
            logger.debug(f"Found {len(slides)} offering slides")
        except Exception as e:
            logger.error(f"Failed to load offering slides from {pptx_path}: {e}", exc_info=True)

        # Only cache default collecte slides
        if pptx_path == self.settings.get_collecte_path(self.base_path):
            self._offerings_cache = slides

        return slides

    # Backwards compatibility alias
    def get_collecte_slides(self, force_refresh: bool = False) -> List[OfferingSlide]:
        """Get slides from default Collecte file. (Backwards compatibility alias)"""
        return self.get_offering_slides(None, force_refresh)

    def _extract_slide_title(self, slide, index: int) -> str:
        """Extract the title from a slide, or generate a default."""
        # Collect all text from shapes for pattern matching
        all_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    all_text.append(text)

        # Try to find "Collecte:" pattern first (more specific for offering slides)
        # Pattern: "Collecte" as whole word, with colon or space, then content up to newline
        # \b ensures word boundary so "Collectedoel" won't match
        collecte_pattern = re.compile(r'(Collecte\b:?\s*[^\r\n\v]*)', re.IGNORECASE)
        for text in all_text:
            match = collecte_pattern.search(text)
            if match:
                title = match.group(1).strip()
                # Make sure we got more than just "Collecte" or "Collecte:"
                if title and not title.rstrip(':').lower().strip() == 'collecte':
                    return title

        # Try to find "Collecte:" pattern in slide notes as fallback
        try:
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                notes_text = notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    match = collecte_pattern.search(notes_text)
                    if match:
                        title = match.group(1).strip()
                        if title and not title.rstrip(':').lower().strip() == 'collecte':
                            return title
        except Exception as e:
            logger.debug(f"Could not access slide notes: {e}")

        # Try to find title shape
        if slide.shapes.title:
            title_text = slide.shapes.title.text.strip()
            if title_text:
                return title_text

        # Try to find any text frame that might be a title
        for text in all_text:
            # Take first line if text is multi-line
            first_line = text.split('\n')[0].strip()
            if first_line and len(first_line) < 100:
                return first_line

        # Default to slide number
        return f"Slide {index + 1}"

    def find_song_by_path(self, relative_path: str) -> Optional[Song]:
        """Find a song by its relative path."""
        songs_path = self.settings.get_songs_path(self.base_path)
        full_path = os.path.join(songs_path, relative_path)
        if os.path.exists(full_path):
            return Song.from_folder(full_path, songs_path)
        return None

    def build_song_tree(self) -> Dict:
        """
        Build a tree structure of the songs folder for display.
        Returns a nested dictionary representing the folder hierarchy.
        """
        songs_path = self.settings.get_songs_path(self.base_path)
        if not os.path.exists(songs_path):
            return {}

        songs = self.scan_songs()
        tree = {}

        for song in songs:
            # Split relative path into parts
            parts = song.relative_path.replace("\\", "/").split("/")

            # Navigate/create tree structure
            current = tree
            for i, part in enumerate(parts[:-1]):  # All but last (the song folder)
                if part not in current:
                    current[part] = {"__children__": {}}
                current = current[part]["__children__"]

            # Add the song at the leaf
            current[parts[-1]] = {"__song__": song}

        return tree
