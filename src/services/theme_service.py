"""Service for loading and saving theme templates as liturgies."""

import os
import shutil
from datetime import datetime
from typing import Optional, List

from pptx import Presentation

from ..models import (
    Liturgy,
    LiturgySection,
    LiturgySlide,
    SectionType,
    Settings,
)
from .pptx_service import PptxService, PptxSection


class ThemeService:
    """Service for working with theme templates."""

    def __init__(self, settings: Settings, base_path: str = "."):
        self.settings = settings
        self.base_path = base_path
        self.pptx_service = PptxService(settings, base_path)

    def load_as_liturgy(self, pptx_path: str, name: Optional[str] = None) -> Liturgy:
        """
        Load a PowerPoint theme file as a Liturgy.
        Parses sections and slides from the PPTX.
        """
        if not os.path.exists(pptx_path):
            raise FileNotFoundError(f"Theme file not found: {pptx_path}")

        # Generate name from filename if not provided
        if not name:
            name = os.path.splitext(os.path.basename(pptx_path))[0]

        # Create liturgy
        liturgy = Liturgy(
            name=name,
            format_version=2,
            theme_source_path=pptx_path,
        )

        # Get sections from the PPTX
        pptx_sections = self.pptx_service.get_sections(pptx_path)

        # Get slides info
        slides_info = self.pptx_service.get_slides_info(pptx_path)

        # If no sections defined, create one section per slide
        if len(pptx_sections) == 1 and pptx_sections[0].name == "All Slides":
            # No sections defined, create individual sections for each slide
            for slide_info in slides_info:
                section = LiturgySection(
                    name=slide_info["title"],
                    section_type=SectionType.REGULAR,
                    source_theme_path=pptx_path,
                )

                # Create slide entry
                slide = LiturgySlide(
                    title=slide_info["title"],
                    slide_index=slide_info["index"],
                    source_path=pptx_path,
                    fields={f.name: f.current_value for f in slide_info["fields"]},
                )
                section.slides.append(slide)

                liturgy.add_section(section)
        else:
            # Use defined sections
            for pptx_section in pptx_sections:
                section = LiturgySection(
                    name=pptx_section.name,
                    section_type=SectionType.REGULAR,
                    source_theme_path=pptx_path,
                )

                # Add slides for this section
                for slide_idx in pptx_section.slide_indices:
                    if slide_idx < len(slides_info):
                        slide_info = slides_info[slide_idx]
                        slide = LiturgySlide(
                            title=slide_info["title"],
                            slide_index=slide_idx,
                            source_path=pptx_path,
                            fields={f.name: f.current_value for f in slide_info["fields"]},
                        )
                        section.slides.append(slide)

                if section.slides:  # Only add non-empty sections
                    liturgy.add_section(section)

        return liturgy

    def save_as_theme(self, liturgy: Liturgy, output_path: str, create_backup: bool = True) -> str:
        """
        Save liturgy as a theme template PPTX.
        Creates timestamped backup if file exists.
        Returns path to the saved file.
        """
        # Create backup if file exists
        if create_backup and os.path.exists(output_path):
            backup_path = self._create_backup(output_path)
            print(f"Created backup: {backup_path}")

        # Merge the liturgy to create the PPTX
        if liturgy.sections:
            temp_path = self.pptx_service.merge_liturgy_v2(liturgy)
        else:
            temp_path = self.pptx_service.merge_liturgy(liturgy)

        # Ensure output directory exists
        os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)

        # Move temp file to output
        shutil.move(temp_path, output_path)

        return output_path

    def _create_backup(self, file_path: str) -> str:
        """
        Create a timestamped backup of a file.
        Format: filename.pptx.20260126T0903
        """
        timestamp = datetime.now().strftime("%Y%m%dT%H%M")
        backup_path = f"{file_path}.{timestamp}"

        # If backup already exists, add seconds
        if os.path.exists(backup_path):
            timestamp = datetime.now().strftime("%Y%m%dT%H%M%S")
            backup_path = f"{file_path}.{timestamp}"

        shutil.copy2(file_path, backup_path)
        return backup_path

    def get_theme_files(self) -> List[str]:
        """Get list of theme files in the themes folder."""
        themes_path = self.settings.get_themes_path(self.base_path)

        if not os.path.exists(themes_path):
            return []

        theme_files = []
        for filename in os.listdir(themes_path):
            if filename.lower().endswith(('.pptx', '.ppt')):
                theme_files.append(os.path.join(themes_path, filename))

        return sorted(theme_files)

    def get_sections_from_theme(self, pptx_path: str) -> List[LiturgySection]:
        """
        Get sections from a theme file without creating a full liturgy.
        Useful for the theme/section picker.
        """
        if not os.path.exists(pptx_path):
            return []

        pptx_sections = self.pptx_service.get_sections(pptx_path)
        slides_info = self.pptx_service.get_slides_info(pptx_path)

        sections = []

        if len(pptx_sections) == 1 and pptx_sections[0].name == "All Slides":
            # No sections defined, create one section per slide
            for slide_info in slides_info:
                section = LiturgySection(
                    name=slide_info["title"],
                    section_type=SectionType.REGULAR,
                    source_theme_path=pptx_path,
                )
                slide = LiturgySlide(
                    title=slide_info["title"],
                    slide_index=slide_info["index"],
                    source_path=pptx_path,
                )
                section.slides.append(slide)
                sections.append(section)
        else:
            for pptx_section in pptx_sections:
                section = LiturgySection(
                    name=pptx_section.name,
                    section_type=SectionType.REGULAR,
                    source_theme_path=pptx_path,
                )

                for slide_idx in pptx_section.slide_indices:
                    if slide_idx < len(slides_info):
                        slide_info = slides_info[slide_idx]
                        slide = LiturgySlide(
                            title=slide_info["title"],
                            slide_index=slide_idx,
                            source_path=pptx_path,
                        )
                        section.slides.append(slide)

                if section.slides:
                    sections.append(section)

        return sections
