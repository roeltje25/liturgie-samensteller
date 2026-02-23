"""Service for scanning existing PPTX files to extract service song lists."""

import enum
import os
import re
from dataclasses import dataclass, field
from datetime import date
from typing import Dict, List, Optional, Set, Tuple

from ..logging_config import get_logger

logger = get_logger("pptx_scanner_service")

# Dutch month names to month numbers
_DUTCH_MONTHS = {
    "januari": 1, "februari": 2, "maart": 3, "april": 4,
    "mei": 5, "juni": 6, "juli": 7, "augustus": 8,
    "september": 9, "oktober": 10, "november": 11, "december": 12,
}

# Regex to detect "Key: value" patterns (not song titles)
_KEY_VALUE_RE = re.compile(r'^\w+\s*:\s', re.IGNORECASE)

# Minimum number of consecutive slides to consider a group a "song"
_MIN_SONG_SLIDES = 2


class SongStatus(enum.Enum):
    """Classification status of a detected song title."""

    CONFIRMED = "confirmed"  # matched a Liederen library folder
    EXCLUDED = "excluded"    # matched an Algemeen item (not a song)
    UNKNOWN = "unknown"      # not found in either


@dataclass
class SongClassification:
    """Classification result for a single detected title."""

    title: str
    status: SongStatus
    library_folder: Optional[str] = None  # relative Liederen path when CONFIRMED


@dataclass
class PptxScanResult:
    """Result of scanning a single PPTX file."""

    filename: str
    filepath: str
    service_date: Optional[date]
    songs: List[str] = field(default_factory=list)
    song_classifications: List[SongClassification] = field(default_factory=list)
    error: Optional[str] = None


class PptxScannerService:
    """Scans existing PPTX presentation files to detect song titles.

    Strategy: groups of 2+ consecutive slides that share an identical short
    text element (≤ 80 chars) are treated as one song.  The shared text is
    cleaned and used as the song title.
    """

    def scan_folder(self, folder_path: str) -> List[PptxScanResult]:
        """Scan all PPTX files in *folder_path* (non-recursive).

        Returns one :class:`PptxScanResult` per ``.pptx`` file found,
        sorted by filename.  Temporary/lock files (starting with ``~``)
        are skipped.
        """
        if not os.path.isdir(folder_path):
            logger.warning("Folder not found: %s", folder_path)
            return []

        filenames = sorted(
            f for f in os.listdir(folder_path)
            if f.lower().endswith(".pptx") and not f.startswith("~")
        )
        return [self.scan_file(os.path.join(folder_path, fn)) for fn in filenames]

    def scan_file(self, filepath: str) -> PptxScanResult:
        """Scan a single PPTX file and return extracted date + songs."""
        filename = os.path.basename(filepath)
        service_date = self.extract_date_from_filename(filename)

        try:
            from pptx import Presentation  # lazy import – keeps startup fast
            prs = Presentation(filepath)
            songs = self.extract_songs_from_pptx(prs)
            return PptxScanResult(
                filename=filename,
                filepath=filepath,
                service_date=service_date,
                songs=songs,
            )
        except Exception as exc:
            logger.error("Error scanning %s: %s", filepath, exc, exc_info=True)
            return PptxScanResult(
                filename=filename,
                filepath=filepath,
                service_date=service_date,
                error=str(exc),
            )

    # ------------------------------------------------------------------
    # Library classification
    # ------------------------------------------------------------------

    def classify_songs(
        self,
        songs: List[str],
        songs_path: str,
        algemeen_path: str,
    ) -> List[SongClassification]:
        """Classify each detected title against the Liederen and Algemeen folders.

        Args:
            songs: List of detected song title strings.
            songs_path: Absolute path to the Liederen root folder.
            algemeen_path: Absolute path to the Algemeen folder.

        Returns:
            A :class:`SongClassification` for every title in *songs*.
        """
        liederen = self._collect_liederen_names(songs_path)
        algemeen = self._collect_algemeen_names(algemeen_path)
        return [
            SongClassification(t, *self._classify_title(t, liederen, algemeen))
            for t in songs
        ]

    def _collect_liederen_names(self, songs_path: str) -> Dict[str, str]:
        """Walk the Liederen root and return {normalized_folder_name: rel_path}.

        Only *leaf* folders — those containing at least one ``.pptx`` or
        ``.pdf`` file — are included.
        """
        result: Dict[str, str] = {}
        if not os.path.isdir(songs_path):
            return result
        for dirpath, _dirnames, filenames in os.walk(songs_path):
            has_song_files = any(
                f.lower().endswith((".pptx", ".pdf")) for f in filenames
            )
            if not has_song_files:
                continue
            folder_name = os.path.basename(dirpath)
            norm = self._normalize(folder_name)
            if norm:
                rel = os.path.relpath(dirpath, songs_path)
                result[norm] = rel
        return result

    def _collect_algemeen_names(self, algemeen_path: str) -> Set[str]:
        """Return normalized names of ``.pptx`` files in the flat Algemeen folder."""
        result: Set[str] = set()
        if not os.path.isdir(algemeen_path):
            return result
        try:
            for entry in os.listdir(algemeen_path):
                if entry.lower().endswith(".pptx"):
                    norm = self._normalize(os.path.splitext(entry)[0])
                    if norm:
                        result.add(norm)
        except OSError as exc:
            logger.warning("Could not list Algemeen folder: %s", exc)
        return result

    def _classify_title(
        self,
        title: str,
        liederen: Dict[str, str],
        algemeen: Set[str],
    ) -> Tuple[SongStatus, Optional[str]]:
        """Classify *title* and return *(status, library_folder)*."""
        norm_title = self._normalize(title)
        title_words = set(norm_title.split())
        significant_title = {w for w in title_words if len(w) > 2}

        # Step 1 — Algemeen check (exclusion)
        for alg_name in algemeen:
            alg_words = set(alg_name.split())
            significant_alg = {w for w in alg_words if len(w) > 2}

            # Forward: ≥80 % of Algemeen's significant words appear in title
            if len(significant_alg) >= 2:
                matches = sum(1 for w in significant_alg if w in title_words)
                if matches / len(significant_alg) >= 0.8:
                    return (SongStatus.EXCLUDED, None)

            # Reverse: all ≥2 significant title words appear in Algemeen name
            if len(significant_title) >= 2:
                if all(w in alg_words for w in significant_title):
                    return (SongStatus.EXCLUDED, None)

        # Step 2 — Liederen check (confirmation)
        best_score = 0.0
        best_path: Optional[str] = None
        for norm_name, rel_path in liederen.items():
            score = self._fuzzy_match(norm_title, norm_name)
            if score > best_score:
                best_score = score
                best_path = rel_path

        if best_score >= 0.75:
            return (SongStatus.CONFIRMED, best_path)

        # Step 3 — Unknown
        return (SongStatus.UNKNOWN, None)

    @staticmethod
    def _normalize(text: str) -> str:
        """Lowercase, strip punctuation (keep alphanumerics + spaces), collapse whitespace."""
        text = text.lower()
        text = re.sub(r"[^\w\s]", " ", text)
        text = re.sub(r"\s+", " ", text).strip()
        return text

    @staticmethod
    def _fuzzy_match(a: str, b: str) -> float:
        """Return a similarity score between 0.0 and 1.0 for strings *a* and *b*."""
        if not a or not b:
            return 0.0
        if a == b:
            return 1.0
        if a in b or b in a:
            return 0.95
        words_a = set(a.split())
        words_b = set(b.split())
        union_w = words_a | words_b
        word_score = (len(words_a & words_b) / len(union_w) * 0.9) if union_w else 0.0
        chars_a = set(a.replace(" ", ""))
        chars_b = set(b.replace(" ", ""))
        union_c = chars_a | chars_b
        char_score = (len(chars_a & chars_b) / len(union_c) * 0.8) if union_c else 0.0
        return max(word_score, char_score)

    # ------------------------------------------------------------------
    # Date extraction
    # ------------------------------------------------------------------

    def extract_date_from_filename(self, filename: str) -> Optional[date]:
        """Try to parse a service date from the filename.

        Supports the following formats (examples from real filenames):

        * ``2025_11_30 hvv dienst.pptx``  → YYYY_MM_DD
        * ``20251116.pptx``               → YYYYMMDD
        * ``14 11 2025.pptx``             → D[D] M[M] YYYY
        * ``29 6 25 viering.pptx``        → D[D] M[M] YY
        * ``17 augustus 2025 …``          → D[D] [dutch-month] YYYY
        """
        name = os.path.splitext(filename)[0]

        # 1. YYYY_MM_DD or YYYY-MM-DD
        m = re.search(r'\b(20\d{2})[_\-](\d{1,2})[_\-](\d{1,2})\b', name)
        if m:
            return self._make_date(int(m.group(1)), int(m.group(2)), int(m.group(3)))

        # 2. YYYYMMDD (eight consecutive digits)
        m = re.search(r'\b(20\d{2})(\d{2})(\d{2})\b', name)
        if m:
            return self._make_date(int(m.group(1)), int(m.group(2)), int(m.group(3)))

        # 3. D[D] M[M] YYYY or D[D] M[M] YY
        m = re.search(r'\b(\d{1,2})\s+(\d{1,2})\s+(20\d{2}|\d{2})\b', name)
        if m:
            day, month, year = int(m.group(1)), int(m.group(2)), int(m.group(3))
            if year < 100:
                year += 2000
            return self._make_date(year, month, day)

        # 4. D[D] [dutch-month-name] YYYY
        month_alt = "|".join(_DUTCH_MONTHS)
        m = re.search(
            rf'\b(\d{{1,2}})\s+({month_alt})\s+(20\d{{2}})\b',
            name,
            re.IGNORECASE,
        )
        if m:
            day = int(m.group(1))
            month = _DUTCH_MONTHS[m.group(2).lower()]
            year = int(m.group(3))
            return self._make_date(year, month, day)

        return None

    @staticmethod
    def _make_date(year: int, month: int, day: int) -> Optional[date]:
        try:
            return date(year, month, day)
        except ValueError:
            return None

    # ------------------------------------------------------------------
    # Song extraction
    # ------------------------------------------------------------------

    def extract_songs_from_pptx(self, prs) -> List[str]:
        """Detect song titles from a merged service presentation.

        Runs two detection phases and merges their results:

        * **Phase 1** — repeated-title: 2+ consecutive slides sharing an
          identical short text element (≤ 80 chars).
        * **Phase 2** — title-then-lyrics: a slide whose ``ph=0`` placeholder
          has short text (≤ 60 chars) is followed by 1–12 lyric slides that
          lack their own title placeholder.

        Returns a deduplicated list of song title strings in order of first
        appearance.
        """
        slides = list(prs.slides)
        n = len(slides)
        if n < _MIN_SONG_SLIDES:
            return []

        groups: List[tuple] = []  # (start, end, text)
        groups.extend(self._find_repeated_title_groups(slides))
        groups.extend(self._find_title_then_lyrics_groups(slides))

        if not groups:
            return []

        # Deduplicate overlapping groups.
        # Sort: by start position, then longer span first so Phase 1's wider
        # span wins over Phase 2 when both detect the same song.
        groups.sort(key=lambda g: (g[0], -(g[1] - g[0])))

        covered: set = set()
        seen_lower: set = set()
        result_with_pos: list = []

        for start, end, text in groups:
            span_set = set(range(start, end + 1))
            if span_set.issubset(covered):
                continue
            title = self._clean_title(text)
            key = title.lower().strip()
            if key in seen_lower:
                continue
            seen_lower.add(key)
            covered |= span_set
            result_with_pos.append((start, title))

        result_with_pos.sort(key=lambda x: x[0])
        return [title for _, title in result_with_pos]

    def _find_repeated_title_groups(self, slides) -> List[tuple]:
        """Phase 1: find groups of 2+ consecutive slides sharing a short title.

        Returns a list of ``(start, end, text)`` tuples (0-based slide indices).
        """
        slide_text_sets = [self._slide_short_texts(slide) for slide in slides]

        text_to_indices: dict = {}
        for i, texts in enumerate(slide_text_sets):
            for t in texts:
                text_to_indices.setdefault(t, []).append(i)

        groups: List[tuple] = []
        for text, indices in text_to_indices.items():
            run_start = indices[0]
            prev = indices[0]
            for idx in indices[1:]:
                if idx == prev + 1:
                    prev = idx
                else:
                    if prev - run_start + 1 >= _MIN_SONG_SLIDES:
                        groups.append((run_start, prev, text))
                    run_start = idx
                    prev = idx
            if prev - run_start + 1 >= _MIN_SONG_SLIDES:
                groups.append((run_start, prev, text))

        return groups

    def _find_title_then_lyrics_groups(self, slides) -> List[tuple]:
        """Phase 2: find title-slide followed by 1–12 lyric slides.

        A *title slide* has a ``ph=0`` placeholder whose text is 3–60 chars
        and passes the title-candidate filter.  A *lyric slide* has no such
        short ``ph=0`` text but contains at least one shape with > 20 chars.

        Returns a list of ``(start, end, text)`` tuples (0-based slide indices).
        """
        n = len(slides)
        groups: List[tuple] = []

        for i in range(n):
            title = self._get_ph0_short_text(slides[i])
            if not title:
                continue

            j = i + 1
            while j < n and (j - i) <= 12:
                if self._is_lyric_slide(slides[j]):
                    j += 1
                else:
                    break

            lyric_count = j - i - 1
            if lyric_count >= 1:
                groups.append((i, j - 1, title))

        return groups

    # ------------------------------------------------------------------
    # Helpers
    # ------------------------------------------------------------------

    def _slide_short_texts(self, slide) -> set:
        """Return the set of candidate title strings from a slide.

        Only includes shape texts whose **total** length is ≤ 80 chars and
        that pass basic "looks like a title" heuristics.
        """
        texts = set()
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            full = shape.text.strip()
            if not full or len(full) > 80:
                continue  # empty or too long → likely lyrics
            if self._is_title_candidate(full):
                texts.add(full)
        return texts

    def _get_ph0_short_text(self, slide) -> Optional[str]:
        """Return the ``ph=0`` placeholder text if it is 3–60 chars and title-like.

        Returns ``None`` when no such placeholder exists or its text is too
        long (indicating it is a lyric rather than a title).
        """
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            try:
                if shape.placeholder_format is None:
                    continue
                if shape.placeholder_format.idx != 0:
                    continue
            except Exception:
                continue
            text = shape.text.strip()
            if 3 <= len(text) <= 60 and self._is_title_candidate(text):
                return text
        return None

    def _is_lyric_slide(self, slide) -> bool:
        """Return True if *slide* looks like a lyric slide (no own title, has content).

        A slide that has its own ``ph=0`` short title text is *not* a lyric
        slide — it would start a new song group.  Otherwise the slide counts
        as a lyric slide when at least one shape has text longer than 20 chars.
        """
        if self._get_ph0_short_text(slide) is not None:
            return False
        for shape in slide.shapes:
            if shape.has_text_frame and len(shape.text.strip()) > 20:
                return True
        return False

    @staticmethod
    def _is_title_candidate(text: str) -> bool:
        """Return True if *text* could plausibly be a song title."""
        if len(text) < 3:
            return False
        # Reject "Key: value" patterns, e.g. "Taal: Engels"
        if _KEY_VALUE_RE.match(text):
            return False
        # Reject bare labels ending with ":", e.g. "FARSI:", "Nederlands:"
        if text.rstrip().endswith(":"):
            return False
        # Reject URLs
        if text.lower().startswith(("http://", "https://", "www.")):
            return False
        # Reject pure numbers
        if re.match(r"^\d+$", text):
            return False
        return True

    @staticmethod
    def _clean_title(text: str) -> str:
        """Return the first non-empty line of *text* (handles \\x0b / \\n)."""
        for line in re.split(r"[\n\x0b]", text):
            line = line.strip()
            if line:
                return line
        return text.strip()
