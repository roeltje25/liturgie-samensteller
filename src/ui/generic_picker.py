"""Dialog for selecting generic items from the Generic (Algemeen) folder."""

from typing import List, Optional

from PyQt6.QtWidgets import (
    QDialog,
    QVBoxLayout,
    QHBoxLayout,
    QListWidget,
    QListWidgetItem,
    QLineEdit,
    QPushButton,
    QLabel,
    QDialogButtonBox,
    QGroupBox,
    QFileDialog,
    QInputDialog,
    QFrame,
)
from PyQt6.QtCore import Qt
from PyQt6.QtGui import QPixmap

from ..models import GenericItem, GenericLiturgyItem
from ..services import PptxService
from ..i18n import tr


class GenericPickerDialog(QDialog):
    """Dialog for selecting a Generic item to add to the liturgy."""

    def __init__(self, items: List[GenericItem], pptx_service: Optional[PptxService] = None, parent=None):
        super().__init__(parent)
        self.items = items
        self.pptx_service = pptx_service
        self._selected_item: Optional[GenericItem] = None
        self._external_path: Optional[str] = None
        self._stub_title: Optional[str] = None

        self._setup_ui()
        self._populate_list()
        self._connect_signals()

    def _setup_ui(self) -> None:
        """Setup the dialog UI."""
        self.setWindowTitle(tr("dialog.generic.title"))
        self.setMinimumSize(500, 500)
        self.resize(600, 600)

        layout = QVBoxLayout(self)

        # Instructions label
        self.instruction_label = QLabel(tr("dialog.generic.select_item"))
        layout.addWidget(self.instruction_label)

        # Search field
        search_layout = QHBoxLayout()
        search_label = QLabel(tr("dialog.generic.search"))
        self.search_input = QLineEdit()
        self.search_input.setPlaceholderText(tr("dialog.generic.search"))
        self.search_input.setClearButtonEnabled(True)
        search_layout.addWidget(search_label)
        search_layout.addWidget(self.search_input)
        layout.addLayout(search_layout)

        # Item list
        self.list_widget = QListWidget()
        self.list_widget.setAlternatingRowColors(True)
        layout.addWidget(self.list_widget)

        # Preview area with thumbnail and info
        preview_layout = QHBoxLayout()

        # Thumbnail
        self.thumbnail_label = QLabel()
        self.thumbnail_label.setFixedSize(120, 90)
        self.thumbnail_label.setFrameStyle(QFrame.Shape.Box | QFrame.Shadow.Sunken)
        self.thumbnail_label.setAlignment(Qt.AlignmentFlag.AlignCenter)
        self.thumbnail_label.setStyleSheet("background-color: #f0f0f0;")
        preview_layout.addWidget(self.thumbnail_label)

        # Info label
        self.info_label = QLabel()
        self.info_label.setWordWrap(True)
        preview_layout.addWidget(self.info_label, 1)

        layout.addLayout(preview_layout)

        # Action buttons group
        action_group = QGroupBox(tr("dialog.generic.actions"))
        action_layout = QHBoxLayout(action_group)

        self.browse_button = QPushButton(tr("button.browse_file"))
        self.browse_button.setToolTip(tr("dialog.generic.browse_tooltip"))
        action_layout.addWidget(self.browse_button)

        self.stub_button = QPushButton(tr("button.create_stub"))
        self.stub_button.setToolTip(tr("dialog.generic.stub_tooltip"))
        action_layout.addWidget(self.stub_button)

        action_layout.addStretch()
        layout.addWidget(action_group)

        # Status label for external/stub selection
        self.status_label = QLabel()
        self.status_label.setStyleSheet("color: blue;")
        layout.addWidget(self.status_label)

        # Button box
        self.button_box = QDialogButtonBox(
            QDialogButtonBox.StandardButton.Ok | QDialogButtonBox.StandardButton.Cancel
        )
        self.button_box.button(QDialogButtonBox.StandardButton.Ok).setText(tr("button.select"))
        self.button_box.button(QDialogButtonBox.StandardButton.Cancel).setText(tr("button.cancel"))
        layout.addWidget(self.button_box)

        # Initially disable OK button
        self.button_box.button(QDialogButtonBox.StandardButton.Ok).setEnabled(False)

    def _populate_list(self) -> None:
        """Populate the list with Generic items."""
        self.list_widget.clear()

        for item in self.items:
            list_item = QListWidgetItem()
            list_item.setText(item.display_name)
            list_item.setData(Qt.ItemDataRole.UserRole, item)
            self.list_widget.addItem(list_item)

    def _connect_signals(self) -> None:
        """Connect widget signals."""
        self.search_input.textChanged.connect(self._on_search_text_changed)
        self.list_widget.itemSelectionChanged.connect(self._on_selection_changed)
        self.list_widget.itemDoubleClicked.connect(self._on_double_click)
        self.browse_button.clicked.connect(self._on_browse_file)
        self.stub_button.clicked.connect(self._on_create_stub)
        self.button_box.accepted.connect(self.accept)
        self.button_box.rejected.connect(self.reject)

    def _on_search_text_changed(self, text: str) -> None:
        """Filter list items based on search text."""
        search_lower = text.lower()
        for i in range(self.list_widget.count()):
            item = self.list_widget.item(i)
            matches = search_lower in item.text().lower()
            item.setHidden(not matches)

    def _on_selection_changed(self) -> None:
        """Handle selection change in list."""
        selected_items = self.list_widget.selectedItems()

        if selected_items:
            item = selected_items[0]
            self._selected_item = item.data(Qt.ItemDataRole.UserRole)
            self._external_path = None
            self._stub_title = None
            self._update_info()
            self.status_label.clear()
        else:
            self._selected_item = None
            self.info_label.setText("")

        self._update_ok_button()

    def _update_info(self) -> None:
        """Update info label with selected item details."""
        if self._selected_item:
            try:
                from pptx import Presentation
                prs = Presentation(self._selected_item.pptx_path)
                slide_count = len(prs.slides)
                self.info_label.setText(tr("dialog.generic.slides", count=slide_count))
            except Exception:
                self.info_label.setText("")

            # Update thumbnail
            self._update_thumbnail(self._selected_item.pptx_path)
        else:
            self.thumbnail_label.setPixmap(QPixmap())

    def _update_thumbnail(self, pptx_path: str) -> None:
        """Update the thumbnail preview."""
        if not self.pptx_service:
            return

        thumb_data = self.pptx_service.get_thumbnail(pptx_path)
        if thumb_data:
            pixmap = QPixmap()
            pixmap.loadFromData(thumb_data)
            if not pixmap.isNull():
                scaled = pixmap.scaled(
                    self.thumbnail_label.size(),
                    Qt.AspectRatioMode.KeepAspectRatio,
                    Qt.TransformationMode.SmoothTransformation
                )
                self.thumbnail_label.setPixmap(scaled)
                return

        self.thumbnail_label.setPixmap(QPixmap())

    def _on_double_click(self, item: QListWidgetItem) -> None:
        """Handle double-click on item."""
        self._selected_item = item.data(Qt.ItemDataRole.UserRole)
        self._external_path = None
        self._stub_title = None
        self.accept()

    def _on_browse_file(self) -> None:
        """Open file dialog to select external PowerPoint file."""
        file_path, _ = QFileDialog.getOpenFileName(
            self,
            tr("dialog.generic.browse_title"),
            "",
            "PowerPoint Files (*.pptx *.ppt);;All Files (*)"
        )

        if file_path:
            self._external_path = file_path
            self._selected_item = None
            self._stub_title = None
            self.list_widget.clearSelection()

            import os
            filename = os.path.basename(file_path)
            self.status_label.setText(tr("dialog.generic.external_selected", filename=filename))
            self.info_label.clear()
            self._update_ok_button()

    def _on_create_stub(self) -> None:
        """Open dialog to create a stub with custom title."""
        title, ok = QInputDialog.getText(
            self,
            tr("dialog.generic.stub_dialog_title"),
            tr("dialog.generic.stub_enter_title"),
        )

        if ok and title.strip():
            self._stub_title = title.strip()
            self._selected_item = None
            self._external_path = None
            self.list_widget.clearSelection()

            self.status_label.setText(tr("dialog.generic.stub_selected", title=self._stub_title))
            self.info_label.clear()
            self._update_ok_button()

    def _update_ok_button(self) -> None:
        """Update OK button enabled state."""
        enabled = (
            self._selected_item is not None
            or self._external_path is not None
            or self._stub_title is not None
        )
        self.button_box.button(QDialogButtonBox.StandardButton.Ok).setEnabled(enabled)

    def get_selected_item(self) -> Optional[GenericLiturgyItem]:
        """Get the selected item as a GenericLiturgyItem."""
        if self._stub_title:
            return GenericLiturgyItem(
                title=self._stub_title,
                is_stub=True,
            )

        if self._external_path:
            import os
            filename = os.path.basename(self._external_path)
            title = os.path.splitext(filename)[0]
            return GenericLiturgyItem(
                title=title,
                source_path=self._external_path,
                pptx_path=self._external_path,
                is_stub=False,
            )

        if self._selected_item:
            return GenericLiturgyItem(
                title=self._selected_item.display_name,
                source_path=self._selected_item.name,
                pptx_path=self._selected_item.pptx_path,
                is_stub=False,
            )

        return None

    # Backwards compatibility alias
    def get_selected_element(self) -> Optional[GenericLiturgyItem]:
        """Get the selected item. (Backwards compatibility alias)"""
        return self.get_selected_item()


# Backwards compatibility alias
AlgemeenPickerDialog = GenericPickerDialog
