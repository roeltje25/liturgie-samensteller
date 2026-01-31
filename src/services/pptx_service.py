"""Service for PowerPoint manipulation and merging."""

import os
import re
import shutil
import subprocess
import tempfile
import time
from typing import List, Optional, Tuple, Dict, Any
from dataclasses import dataclass, field

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor

from ..logging_config import get_logger

logger = get_logger("pptx_service")

from ..models import (
    Settings,
    Liturgy,
    LiturgyItem,
    SongLiturgyItem,
    GenericLiturgyItem,
    OfferingLiturgyItem,
    ItemType,
    LiturgySection,
    LiturgySlide,
    SectionType,
)


@dataclass
class PptxSection:
    """Represents a section in a PowerPoint file."""

    name: str
    slide_indices: List[int] = field(default_factory=list)


@dataclass
class SlideField:
    """Represents a fillable field in a slide."""

    name: str
    field_type: str  # "placeholder" or "text_pattern"
    placeholder_idx: Optional[int] = None  # For placeholder fields
    current_value: str = ""


class PptxService:
    """Service for creating and manipulating PowerPoint presentations."""

    def __init__(self, settings: Settings, base_path: str = "."):
        self.settings = settings
        self.base_path = base_path

    def merge_liturgy(self, liturgy: Liturgy) -> str:
        """
        Create a merged presentation from all liturgy items.
        Returns path to the merged file.

        Tries merge methods in order:
        1. VBScript executed by PowerPoint (Windows, best formatting preservation)
        2. python-pptx fallback (may lose some formatting)
        """
        # Collect all source files and slide indices to copy
        slides_to_copy = []  # List of (pptx_path, slide_indices)

        for item in liturgy.items:
            slide_info = self._get_slides_for_item(item)
            if slide_info:
                slides_to_copy.append(slide_info)

        if not slides_to_copy:
            # No slides to copy, create empty presentation
            logger.info("merge_liturgy: No slides to copy, creating empty presentation")
            try:
                prs = Presentation()
                prs.slide_width = Inches(13.333)
                prs.slide_height = Inches(7.5)
                temp_file = tempfile.NamedTemporaryFile(suffix='.pptx', delete=False)
                temp_path = temp_file.name
                temp_file.close()
                prs.save(temp_path)
                logger.debug(f"Empty presentation saved to: {temp_path}")
                return temp_path
            except Exception as e:
                logger.error(f"Failed to create empty presentation: {e}", exc_info=True)
                raise

        # Try VBScript method (Windows with PowerPoint installed)
        if os.name == 'nt':
            try:
                return self._merge_with_vbscript(slides_to_copy)
            except Exception as e:
                print(f"VBScript merge failed: {e}, falling back to python-pptx")

        # Fallback to python-pptx method
        return self._merge_with_pptx(slides_to_copy)

    def _get_slides_for_item(self, item: LiturgyItem) -> Optional[Tuple[str, List[int]]]:
        """Get the source file and slide indices for a liturgy item."""
        if item.is_stub:
            # Create a stub presentation
            stub_path = self._create_stub_presentation(item.title)
            return (stub_path, [0])

        if item.item_type == ItemType.SONG:
            song_item: SongLiturgyItem = item
            if song_item.pptx_path and os.path.exists(song_item.pptx_path):
                slide_count = self.get_slide_count(song_item.pptx_path)
                return (song_item.pptx_path, list(range(slide_count)))

        elif item.item_type == ItemType.GENERIC:
            generic_item: GenericLiturgyItem = item
            if generic_item.pptx_path and os.path.exists(generic_item.pptx_path):
                slide_count = self.get_slide_count(generic_item.pptx_path)
                return (generic_item.pptx_path, list(range(slide_count)))

        elif item.item_type == ItemType.OFFERING:
            offering_item: OfferingLiturgyItem = item
            pptx_path = offering_item.pptx_path or self.settings.get_collecte_path(self.base_path)
            if os.path.exists(pptx_path):
                return (pptx_path, [offering_item.slide_index])

        return None

    def _create_stub_presentation(self, title: str) -> str:
        """Create a stub presentation with a single slide."""
        template_path = self.settings.get_stub_template_path(self.base_path)

        if template_path and os.path.exists(template_path):
            # Copy template and modify title
            temp_path = tempfile.mktemp(suffix='.pptx')
            shutil.copy(template_path, temp_path)

            # Try to set title
            try:
                prs = Presentation(temp_path)
                if len(prs.slides) > 0 and prs.slides[0].shapes.title:
                    prs.slides[0].shapes.title.text = title
                prs.save(temp_path)
            except Exception:
                pass

            return temp_path
        else:
            # Create simple stub
            logger.info(f"Creating simple stub presentation for: {title}")
            try:
                prs = Presentation()
                prs.slide_width = Inches(13.333)
                prs.slide_height = Inches(7.5)

                blank_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(blank_layout)

                left = Inches(1)
                top = Inches(3)
                width = Inches(11.333)
                height = Inches(1.5)

                textbox = slide.shapes.add_textbox(left, top, width, height)
                tf = textbox.text_frame
                p = tf.paragraphs[0]
                p.text = title
                p.font.size = Pt(44)
                p.font.bold = True
                p.alignment = PP_ALIGN.CENTER

                temp_file = tempfile.NamedTemporaryFile(suffix='.pptx', delete=False)
                temp_path = temp_file.name
                temp_file.close()
                prs.save(temp_path)
                logger.debug(f"Stub presentation saved to: {temp_path}")
                return temp_path
            except Exception as e:
                logger.error(f"Failed to create stub presentation: {e}", exc_info=True)
                raise

    def _merge_with_vbscript(self, slides_to_copy: List[Tuple[str, List[int]]]) -> str:
        """
        Merge presentations by generating and executing a VBScript in PowerPoint.
        This approach lets PowerPoint handle the merge natively, preserving all formatting.
        Uses InsertFromFile to properly handle Designer-created slides and special elements.
        """
        temp_path = tempfile.mktemp(suffix='.pptx')
        abs_temp = os.path.abspath(temp_path).replace("/", "\\")

        # Build the VBScript
        vbs_lines = [
            'Option Explicit',
            '',
            'Const ppSaveAsOpenXMLPresentation = 24',
            'Const ppWindowMinimized = 2',
            '',
            'Dim pptApp, targetPres, sourcePres, baseDesign',
            'Dim fso, slideIndex, insertedSlide',
            '',
            'On Error Resume Next',
            '',
            '\'Create PowerPoint application',
            'Set pptApp = CreateObject("PowerPoint.Application")',
            'If Err.Number <> 0 Then',
            '    WScript.Echo "Error: Could not start PowerPoint. " & Err.Description',
            '    WScript.Quit 1',
            'End If',
            'On Error GoTo 0',
            '',
            '\'Make PowerPoint visible but minimized',
            'pptApp.Visible = True',
            'pptApp.WindowState = ppWindowMinimized',
            '',
            '\'Create file system object for file existence checks',
            'Set fso = CreateObject("Scripting.FileSystemObject")',
            '',
            '\'Create a new presentation as target',
            'Set targetPres = pptApp.Presentations.Add(True)',
            '',
            '\'Track where to insert next slide',
            'slideIndex = 0',
            '',
            '\'Give PowerPoint time to initialize',
            'WScript.Sleep 500',
            '',
        ]

        # Add slide insertion commands for each source file
        for source_path, slide_indices in slides_to_copy:
            if not os.path.exists(source_path):
                continue

            abs_source = os.path.abspath(source_path).replace("/", "\\")

            # Escape backslashes for VBScript string
            escaped_source = abs_source.replace("\\", "\\\\")

            vbs_lines.append(f'\'Insert slides from: {os.path.basename(source_path)}')
            vbs_lines.append(f'If fso.FileExists("{escaped_source}") Then')

            # For each slide index, use InsertFromFile
            for slide_idx in slide_indices:
                # VBScript/PowerPoint uses 1-based indices
                ppt_idx = slide_idx + 1
                vbs_lines.append(f'    \'Insert slide {ppt_idx}')
                vbs_lines.append(f'    On Error Resume Next')
                # InsertFromFile(FileName, Index, SlideStart, SlideEnd)
                # Index = where to insert (position after which to insert, 0 = at beginning)
                # SlideStart/SlideEnd = which slides from source (1-based)
                vbs_lines.append(f'    targetPres.Slides.InsertFromFile "{escaped_source}", slideIndex, {ppt_idx}, {ppt_idx}')
                vbs_lines.append(f'    If Err.Number = 0 Then')
                vbs_lines.append(f'        slideIndex = slideIndex + 1')
                vbs_lines.append(f'    End If')
                vbs_lines.append(f'    On Error GoTo 0')
                vbs_lines.append(f'    WScript.Sleep 100')
                vbs_lines.append(f'')

            vbs_lines.append('End If')
            vbs_lines.append('')

        # Escape the output path
        escaped_output = abs_temp.replace("\\", "\\\\")

        # Add save and cleanup code
        vbs_lines.extend([
            '\'Remove the initial blank slide if we added slides',
            'If targetPres.Slides.Count > 1 Then',
            '    On Error Resume Next',
            '    Dim firstSlideShapeCount',
            '    firstSlideShapeCount = targetPres.Slides(1).Shapes.Count',
            '    If firstSlideShapeCount = 0 Then',
            '        targetPres.Slides(1).Delete',
            '    End If',
            '    On Error GoTo 0',
            'End If',
            '',
            '\'Save as PPTX',
            f'targetPres.SaveAs "{escaped_output}", ppSaveAsOpenXMLPresentation',
            '',
            '\'Close and cleanup',
            'targetPres.Close',
            'pptApp.Quit',
            '',
            'Set targetPres = Nothing',
            'Set sourcePres = Nothing',
            'Set pptApp = Nothing',
            'Set fso = Nothing',
            '',
            'WScript.Echo "SUCCESS"',
            'WScript.Quit 0',
        ])

        # Write VBScript to temp file
        vbs_content = '\r\n'.join(vbs_lines)
        vbs_path = tempfile.mktemp(suffix='.vbs')

        with open(vbs_path, 'w', encoding='utf-8') as f:
            f.write(vbs_content)

        try:
            # Execute the VBScript using cscript (console mode for better error handling)
            result = subprocess.run(
                ['cscript', '//Nologo', vbs_path],
                capture_output=True,
                text=True,
                timeout=300  # 5 minute timeout
            )

            # Check for success
            if result.returncode != 0:
                error_msg = result.stderr.strip() or result.stdout.strip() or "Unknown error"
                raise RuntimeError(f"VBScript execution failed: {error_msg}")

            if "SUCCESS" not in result.stdout:
                raise RuntimeError(f"VBScript did not complete successfully: {result.stdout}")

            # Verify output file exists
            if not os.path.exists(temp_path):
                raise RuntimeError("VBScript completed but output file was not created")

            return temp_path

        finally:
            # Clean up VBScript file
            try:
                os.remove(vbs_path)
            except Exception:
                pass

    def _duplicate_slide(self, pres, source_slide):
        """
        Duplicate a slide within a presentation, preserving all content.
        """
        import copy
        from lxml import etree

        # Find the blank layout (one with fewest placeholders)
        layout_items_count = [len(layout.placeholders) for layout in pres.slide_layouts]
        min_items = min(layout_items_count)
        blank_layout_id = layout_items_count.index(min_items)
        blank_layout = pres.slide_layouts[blank_layout_id]

        # Add new slide
        new_slide = pres.slides.add_slide(blank_layout)

        # Copy shapes via XML
        for shape in source_slide.shapes:
            el = shape.element
            new_el = copy.deepcopy(el)
            new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')

        # Copy relationships (images, etc.)
        for rel in source_slide.part.rels.values():
            # Skip notes and slide layout relationships
            if "notesSlide" in rel.reltype or "slideLayout" in rel.reltype:
                continue

            try:
                if rel.is_external:
                    new_slide.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
                else:
                    new_slide.part.rels.get_or_add(rel.reltype, rel.target_part)
            except Exception:
                pass

        return new_slide

    def _copy_slide_from_source(self, target_pres, source_pres, slide_idx):
        """
        Copy a slide from source presentation to target presentation.
        Handles images, shapes, and formatting.
        """
        import copy
        from io import BytesIO
        from pptx.util import Emu
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        if slide_idx >= len(source_pres.slides):
            return None

        source_slide = source_pres.slides[slide_idx]

        # Find blank layout in target
        layout_items_count = [len(layout.placeholders) for layout in target_pres.slide_layouts]
        min_items = min(layout_items_count)
        blank_layout_id = layout_items_count.index(min_items)
        blank_layout = target_pres.slide_layouts[blank_layout_id]

        new_slide = target_pres.slides.add_slide(blank_layout)

        # Copy each shape
        for shape in source_slide.shapes:
            try:
                self._copy_shape_to_slide(shape, new_slide, source_slide)
            except Exception as e:
                print(f"Failed to copy shape: {e}")

        # Copy slide background if present
        try:
            self._copy_slide_background(source_slide, new_slide)
        except Exception as e:
            logger.debug(f"Failed to copy background: {e}")

        return new_slide

    def _copy_slide_background(self, source_slide, target_slide) -> None:
        """Copy the background from source slide to target slide."""
        from pptx.enum.dml import MSO_FILL_TYPE
        from io import BytesIO

        source_bg = source_slide.background
        target_bg = target_slide.background

        # Check if source has a background fill
        if source_bg.fill.type is None:
            return

        fill_type = source_bg.fill.type

        if fill_type == MSO_FILL_TYPE.SOLID:
            # Solid color background
            try:
                target_bg.fill.solid()
                # Try to get the actual RGB color
                fore_color = source_bg.fill.fore_color
                if fore_color.type == MSO_THEME_COLOR.NOT_THEME_COLOR:
                    # Direct RGB color
                    target_bg.fill.fore_color.rgb = fore_color.rgb
                else:
                    # Theme color - try to get the actual color value
                    try:
                        target_bg.fill.fore_color.rgb = fore_color.rgb
                    except Exception:
                        # Fall back to theme color reference
                        target_bg.fill.fore_color.theme_color = fore_color.theme_color
            except Exception as e:
                logger.debug(f"Failed to copy solid background: {e}")

        elif fill_type == MSO_FILL_TYPE.PICTURE:
            # Picture/image background
            try:
                image_blob = source_bg.fill._fill._pic.blob
                image_stream = BytesIO(image_blob)
                target_bg.fill.background()
                # python-pptx doesn't easily support setting picture backgrounds
                # This is a limitation of the library
                logger.debug("Picture backgrounds not fully supported in python-pptx")
            except Exception as e:
                logger.debug(f"Failed to copy picture background: {e}")

        elif fill_type == MSO_FILL_TYPE.GRADIENT:
            # Gradient background - complex to copy
            try:
                # Try to copy at least the start and end colors
                target_bg.fill.gradient()
                logger.debug("Gradient backgrounds partially supported")
            except Exception as e:
                logger.debug(f"Failed to copy gradient background: {e}")

    def _copy_shape_to_slide(self, shape, target_slide, source_slide):
        """Copy a single shape to target slide."""
        from io import BytesIO
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from pptx.util import Pt
        import copy

        # Handle pictures
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                image_blob = shape.image.blob
                image_stream = BytesIO(image_blob)
                target_slide.shapes.add_picture(
                    image_stream,
                    shape.left, shape.top,
                    shape.width, shape.height
                )
            except Exception:
                pass
            return

        # Handle groups
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            # For groups, copy the XML element directly
            try:
                new_el = copy.deepcopy(shape.element)
                target_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
                # Also need to copy any image relationships within the group
                self._copy_group_images(shape, target_slide, source_slide)
            except Exception:
                pass
            return

        # Handle tables
        if shape.has_table:
            try:
                table = shape.table
                rows = len(table.rows)
                cols = len(table.columns)

                # Get column widths
                col_widths = [col.width for col in table.columns]

                new_table = target_slide.shapes.add_table(
                    rows, cols,
                    shape.left, shape.top,
                    shape.width, shape.height
                ).table

                # Copy cell contents
                for row_idx, row in enumerate(table.rows):
                    for col_idx, cell in enumerate(row.cells):
                        target_cell = new_table.cell(row_idx, col_idx)
                        target_cell.text = cell.text
                        # Copy paragraph formatting
                        if cell.text_frame and target_cell.text_frame:
                            for p_idx, para in enumerate(cell.text_frame.paragraphs):
                                if p_idx < len(target_cell.text_frame.paragraphs):
                                    target_para = target_cell.text_frame.paragraphs[p_idx]
                                    try:
                                        target_para.alignment = para.alignment
                                    except Exception:
                                        pass
            except Exception:
                pass
            return

        # Handle text boxes and other shapes with text
        if shape.has_text_frame:
            try:
                # Create textbox
                new_shape = target_slide.shapes.add_textbox(
                    shape.left, shape.top,
                    shape.width, shape.height
                )

                # Copy text frame content
                target_tf = new_shape.text_frame
                source_tf = shape.text_frame

                # Copy text frame properties
                target_tf.word_wrap = source_tf.word_wrap

                for p_idx, source_para in enumerate(source_tf.paragraphs):
                    if p_idx == 0:
                        target_para = target_tf.paragraphs[0]
                    else:
                        target_para = target_tf.add_paragraph()

                    # Copy runs (text with formatting)
                    for r_idx, run in enumerate(source_para.runs):
                        if r_idx == 0 and p_idx == 0:
                            target_run = target_para.runs[0] if target_para.runs else target_para.add_run()
                        else:
                            target_run = target_para.add_run()

                        target_run.text = run.text

                        # Copy font properties
                        try:
                            if run.font.size:
                                target_run.font.size = run.font.size
                            if run.font.bold is not None:
                                target_run.font.bold = run.font.bold
                            if run.font.italic is not None:
                                target_run.font.italic = run.font.italic
                            if run.font.name:
                                target_run.font.name = run.font.name
                            # Copy font color - handle both RGB and theme colors
                            if run.font.color:
                                try:
                                    # Try direct RGB first
                                    if run.font.color.rgb:
                                        target_run.font.color.rgb = run.font.color.rgb
                                except Exception:
                                    try:
                                        # Try to get theme color and convert to RGB
                                        if run.font.color.theme_color is not None:
                                            # Get the actual RGB value from the theme
                                            rgb_val = run.font.color._color.rgb
                                            if rgb_val:
                                                target_run.font.color.rgb = rgb_val
                                    except Exception:
                                        pass
                        except Exception:
                            pass

                    # If no runs, copy paragraph text directly
                    if not source_para.runs:
                        target_para.text = source_para.text

                    # Copy paragraph formatting
                    try:
                        target_para.alignment = source_para.alignment
                        if source_para.level is not None:
                            target_para.level = source_para.level
                    except Exception:
                        pass

            except Exception:
                pass
            return

        # For other shapes, try deep copy of element
        try:
            new_el = copy.deepcopy(shape.element)
            target_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
        except Exception:
            pass

    def _copy_group_images(self, group_shape, target_slide, source_slide):
        """Copy image relationships from a group shape."""
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from io import BytesIO

        try:
            for shape in group_shape.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    # The image was copied via XML but relationship may be missing
                    # Try to add the relationship
                    try:
                        image_blob = shape.image.blob
                        # Add to target's relationships
                        target_slide.part.get_or_add_image_part(BytesIO(image_blob))
                    except Exception:
                        pass
                elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    self._copy_group_images(shape, target_slide, source_slide)
        except Exception:
            pass

    def _merge_with_pptx(self, slides_to_copy: List[Tuple[str, List[int]]]) -> str:
        """
        Fallback merge using python-pptx.
        Uses the first source as base and copies slides from other sources.
        """
        logger.info(f"Merging {len(slides_to_copy)} presentations using python-pptx")
        # Use first presentation as base - this preserves its theme/master
        first_path, first_indices = slides_to_copy[0]
        try:
            prs = Presentation(first_path)
            logger.debug(f"Loaded base presentation: {first_path}")
        except Exception as e:
            logger.error(f"Failed to load base presentation {first_path}: {e}", exc_info=True)
            raise

        # Remove slides we don't want from the first presentation
        all_indices = set(range(len(prs.slides)))
        indices_to_remove = all_indices - set(first_indices)

        # Remove in reverse order to maintain indices
        for idx in sorted(indices_to_remove, reverse=True):
            rId = prs.slides._sldIdLst[idx].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[idx]

        # Add slides from other presentations
        for source_path, slide_indices in slides_to_copy[1:]:
            if not os.path.exists(source_path):
                logger.warning(f"Source presentation not found, skipping: {source_path}")
                continue

            try:
                source_prs = Presentation(source_path)
                logger.debug(f"Loading source presentation: {source_path}")
                for slide_idx in slide_indices:
                    self._copy_slide_from_source(prs, source_prs, slide_idx)
            except Exception as e:
                logger.error(f"Failed to load/copy from {source_path}: {e}", exc_info=True)
                continue

        temp_file = tempfile.NamedTemporaryFile(suffix='.pptx', delete=False)
        temp_path = temp_file.name
        temp_file.close()
        prs.save(temp_path)
        logger.info(f"Merged presentation saved to: {temp_path}")
        return temp_path

    def save_presentation(self, temp_path: str, output_path: str) -> None:
        """Save presentation from temp path to output path."""
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        shutil.move(temp_path, output_path)

    def get_slide_count(self, pptx_path: str) -> int:
        """Get the number of slides in a presentation."""
        if not os.path.exists(pptx_path):
            return 0
        try:
            prs = Presentation(pptx_path)
            return len(prs.slides)
        except Exception:
            return 0

    def get_thumbnail(self, pptx_path: str) -> Optional[bytes]:
        """Extract the thumbnail image from a PowerPoint file.

        Returns the image bytes or None if not available.
        """
        if not os.path.exists(pptx_path):
            return None
        try:
            import zipfile
            with zipfile.ZipFile(pptx_path, 'r') as zf:
                # Try common thumbnail locations
                for thumb_path in ['docProps/thumbnail.jpeg', 'docProps/thumbnail.png']:
                    if thumb_path in zf.namelist():
                        return zf.read(thumb_path)
        except Exception:
            pass
        return None

    def get_slide_thumbnail(self, pptx_path: str, slide_index: int, width: int = 240) -> Optional[bytes]:
        """Get a thumbnail image for a specific slide.

        Uses PowerPoint COM automation on Windows to export the slide as an image.
        Results are cached to avoid repeated exports.

        Args:
            pptx_path: Path to the PowerPoint file
            slide_index: 0-based index of the slide
            width: Desired width of the thumbnail in pixels

        Returns the image bytes or None if not available.
        """
        if not os.path.exists(pptx_path):
            return None

        # Create cache directory
        cache_dir = os.path.join(tempfile.gettempdir(), "pptx_thumbnails")
        os.makedirs(cache_dir, exist_ok=True)

        # Generate cache key based on file path, modification time, and slide index
        file_stat = os.stat(pptx_path)
        cache_key = f"{os.path.basename(pptx_path)}_{file_stat.st_mtime}_{slide_index}_{width}"
        # Create a safe filename from the cache key
        import hashlib
        cache_filename = hashlib.md5(cache_key.encode()).hexdigest() + ".png"
        cache_path = os.path.join(cache_dir, cache_filename)

        # Check if cached thumbnail exists
        if os.path.exists(cache_path):
            try:
                with open(cache_path, 'rb') as f:
                    return f.read()
            except Exception:
                pass

        # Try to generate thumbnail using PowerPoint COM
        try:
            thumb_data = self._export_slide_with_com(pptx_path, slide_index, cache_path, width)
            if thumb_data:
                return thumb_data
        except Exception:
            pass

        # Fallback: return the presentation thumbnail (first slide only)
        if slide_index == 0:
            return self.get_thumbnail(pptx_path)

        return None

    def _export_slide_with_com(self, pptx_path: str, slide_index: int,
                                output_path: str, width: int) -> Optional[bytes]:
        """Export a slide as an image using PowerPoint COM automation.

        Args:
            pptx_path: Path to the PowerPoint file
            slide_index: 0-based index of the slide
            output_path: Path to save the exported image
            width: Desired width of the image

        Returns the image bytes or None if export failed.
        """
        try:
            import win32com.client
            import pythoncom
        except ImportError:
            logger.debug("win32com not available for slide thumbnail export")
            return None

        powerpoint = None
        presentation = None
        com_initialized = False

        try:
            # Initialize COM in this thread
            pythoncom.CoInitialize()
            com_initialized = True
            logger.debug(f"Exporting slide {slide_index} from {pptx_path}")

            # Create PowerPoint application
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")

            # Open presentation (read-only, no window)
            presentation = powerpoint.Presentations.Open(
                os.path.abspath(pptx_path),
                ReadOnly=True,
                Untitled=False,
                WithWindow=False
            )

            # Get the slide (1-based index in COM)
            slide_number = slide_index + 1
            if slide_number > presentation.Slides.Count:
                logger.warning(f"Slide index {slide_index} out of range (max {presentation.Slides.Count - 1})")
                return None

            slide = presentation.Slides(slide_number)

            # Calculate height based on aspect ratio (assume 16:9)
            height = int(width * 9 / 16)

            # Export slide as PNG
            slide.Export(output_path, "PNG", width, height)

            # Read the exported image
            if os.path.exists(output_path):
                with open(output_path, 'rb') as f:
                    return f.read()

        except Exception as e:
            logger.error(f"Failed to export slide thumbnail: {pptx_path}[{slide_index}]: {e}", exc_info=True)
        finally:
            # Clean up COM objects
            try:
                if presentation:
                    presentation.Close()
            except Exception as e:
                logger.debug(f"Error closing presentation: {e}")
            try:
                if powerpoint:
                    powerpoint.Quit()
            except Exception as e:
                logger.debug(f"Error quitting PowerPoint: {e}")
            if com_initialized:
                try:
                    pythoncom.CoUninitialize()
                except Exception as e:
                    logger.debug(f"Error uninitializing COM: {e}")

        return None

    def get_slide_thumbnail_text(self, pptx_path: str, slide_index: int) -> str:
        """Get a text representation of a slide (title + first paragraph)."""
        if not os.path.exists(pptx_path):
            return ""
        try:
            prs = Presentation(pptx_path)
            if slide_index < len(prs.slides):
                slide = prs.slides[slide_index]
                texts = []
                if slide.shapes.title:
                    texts.append(self._clean_title(slide.shapes.title.text))
                for shape in slide.shapes:
                    if shape.has_text_frame and shape != slide.shapes.title:
                        text = self._clean_title(shape.text_frame.text)
                        if text:
                            texts.append(text[:100])
                            break
                return "\n".join(texts)
        except Exception:
            pass
        return ""

    def get_sections(self, pptx_path: str) -> List[PptxSection]:
        """
        Read PowerPoint sections from a presentation.
        Returns list of PptxSection objects with section names and slide indices.
        """
        if not os.path.exists(pptx_path):
            return []

        try:
            import zipfile
            from lxml import etree

            prs = Presentation(pptx_path)
            sections = []
            total_slides = len(prs.slides)

            # Read presentation.xml directly from the ZIP file
            # (prs.part.element doesn't always work correctly)
            with zipfile.ZipFile(pptx_path, 'r') as z:
                with z.open('ppt/presentation.xml') as f:
                    tree = etree.parse(f)
                    root = tree.getroot()

            namespaces = {
                'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
                'p14': 'http://schemas.microsoft.com/office/powerpoint/2010/main',
                'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
            }

            # Build a map of slide IDs to slide indices
            slide_id_to_index = {}
            sld_id_lst = root.find('.//p:sldIdLst', namespaces)
            if sld_id_lst is not None:
                for idx, sld_id in enumerate(sld_id_lst.findall('p:sldId', namespaces)):
                    slide_id = sld_id.get('id')
                    if slide_id:
                        slide_id_to_index[slide_id] = idx

            # Look for section list element
            section_list = root.find('.//p14:sectionLst', namespaces)

            if section_list is not None:
                for section_elem in section_list.findall('p14:section', namespaces):
                    section_name = section_elem.get('name', 'Untitled')
                    # Clean section name of newlines
                    section_name = self._clean_title(section_name)

                    # Get slide IDs in this section
                    slide_id_list = section_elem.find('p14:sldIdLst', namespaces)
                    slide_indices = []

                    if slide_id_list is not None:
                        for sld_id in slide_id_list.findall('p14:sldId', namespaces):
                            # The id attribute references the slide ID
                            ref_id = sld_id.get('id')
                            if ref_id and ref_id in slide_id_to_index:
                                slide_indices.append(slide_id_to_index[ref_id])

                    # Sort indices to ensure proper order
                    slide_indices.sort()

                    sections.append(PptxSection(
                        name=section_name,
                        slide_indices=slide_indices,
                    ))

                if sections:
                    return sections

            # Fallback: if no sections found, treat whole presentation as one section
            sections.append(PptxSection(
                name="All Slides",
                slide_indices=list(range(total_slides)),
            ))

            return sections

        except Exception as e:
            logger.error(f"Error reading sections from {pptx_path}: {e}", exc_info=True)
            return []

    def _clean_title(self, title: str) -> str:
        """Clean a title by removing newlines and excess whitespace."""
        if not title:
            return title
        # Replace various newline characters with spaces
        cleaned = title.replace('\r\n', ' ').replace('\n', ' ').replace('\r', ' ')
        # Replace vertical tab and other control characters
        cleaned = ''.join(c if c.isprintable() or c == ' ' else ' ' for c in cleaned)
        # Collapse multiple spaces
        cleaned = ' '.join(cleaned.split())
        return cleaned.strip()

    def get_slides_info(self, pptx_path: str) -> List[Dict[str, Any]]:
        """
        Get information about all slides in a presentation.
        Returns list of dicts with slide title, index, and extracted fields.
        """
        if not os.path.exists(pptx_path):
            return []

        try:
            prs = Presentation(pptx_path)
            slides_info = []

            for idx, slide in enumerate(prs.slides):
                # Get slide title
                title = ""
                if slide.shapes.title:
                    title = self._clean_title(slide.shapes.title.text)

                # Extract fields
                fields = self.extract_fields_from_slide(slide)

                slides_info.append({
                    "index": idx,
                    "title": title or f"Slide {idx + 1}",
                    "fields": fields,
                })

            return slides_info

        except Exception as e:
            print(f"Error getting slides info from {pptx_path}: {e}")
            return []

    def extract_fields_from_slide(self, slide) -> List[SlideField]:
        """
        Extract fillable fields from a slide.
        Looks for:
        1. Native placeholders (Title, Subtitle, Body, etc.)
        2. Text patterns like {FIELD_NAME}
        """
        fields = []
        # Match {FieldName} patterns - allows mixed case, letters, numbers, underscores
        field_pattern = re.compile(r'\{([A-Za-z_][A-Za-z0-9_]*)\}')

        # Map placeholder type names to user-friendly names
        placeholder_name_map = {
            'TITLE': 'TITLE',
            'CENTER_TITLE': 'TITLE',
            'SUBTITLE': 'SUBTITLE',
            'BODY': 'BODY',
            'OBJECT': 'CONTENT',
            'CHART': 'CHART',
            'TABLE': 'TABLE',
            'PICTURE': 'PICTURE',
            'MEDIA_CLIP': 'MEDIA',
            'ORG_CHART': 'ORG_CHART',
            'DATE': 'DATE',
            'FOOTER': 'FOOTER',
            'SLIDE_NUMBER': 'SLIDE_NUMBER',
            'HEADER': 'HEADER',
        }

        # Track placeholder counts to handle duplicates
        placeholder_counts = {}

        # 1. Check placeholders
        try:
            for placeholder in slide.placeholders:
                if placeholder.has_text_frame:
                    text = placeholder.text_frame.text
                    ph_idx = placeholder.placeholder_format.idx

                    # Determine field name from placeholder type
                    field_name = f"PLACEHOLDER_{ph_idx}"
                    try:
                        ph_type = placeholder.placeholder_format.type
                        if ph_type:
                            # Extract the base type name (e.g., "BODY" from "BODY (2)")
                            type_str = str(ph_type)
                            base_name = type_str.split(' ')[0].split('(')[0].upper()

                            # Map to user-friendly name
                            field_name = placeholder_name_map.get(base_name, base_name)
                    except Exception:
                        pass

                    # Handle duplicate field names by appending index
                    if field_name in placeholder_counts:
                        placeholder_counts[field_name] += 1
                        field_name = f"{field_name}_{placeholder_counts[field_name]}"
                    else:
                        placeholder_counts[field_name] = 1

                    fields.append(SlideField(
                        name=field_name,
                        field_type="placeholder",
                        placeholder_idx=ph_idx,
                        current_value=text,
                    ))
        except Exception:
            pass

        # 2. Check for {FIELD_NAME} patterns in all text
        try:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            matches = field_pattern.findall(run.text)
                            for match in matches:
                                # Avoid duplicates
                                if not any(f.name == match and f.field_type == "text_pattern" for f in fields):
                                    fields.append(SlideField(
                                        name=match,
                                        field_type="text_pattern",
                                        current_value=f"{{{match}}}",
                                    ))
        except Exception:
            pass

        return fields

    def extract_fields(self, pptx_path: str, slide_index: int = None) -> List[SlideField]:
        """
        Extract fillable fields from a presentation.
        If slide_index is provided, extract from that slide only.
        Otherwise, extract from all slides.
        """
        if not os.path.exists(pptx_path):
            return []

        try:
            prs = Presentation(pptx_path)
            all_fields = []

            if slide_index is not None:
                if 0 <= slide_index < len(prs.slides):
                    all_fields = self.extract_fields_from_slide(prs.slides[slide_index])
            else:
                for slide in prs.slides:
                    all_fields.extend(self.extract_fields_from_slide(slide))

            return all_fields

        except Exception as e:
            print(f"Error extracting fields from {pptx_path}: {e}")
            return []

    def fill_slide_fields(self, slide, fields: Dict[str, str]) -> None:
        """
        Fill fields in a slide with provided values.
        Handles both placeholders and {FIELD_NAME} text patterns.
        """
        if not fields:
            return

        # Match {FieldName} patterns - allows mixed case
        field_pattern = re.compile(r'\{([A-Za-z_][A-Za-z0-9_]*)\}')

        # Map placeholder type names to user-friendly names (same as in extract)
        placeholder_name_map = {
            'TITLE': 'TITLE',
            'CENTER_TITLE': 'TITLE',
            'SUBTITLE': 'SUBTITLE',
            'BODY': 'BODY',
            'OBJECT': 'CONTENT',
            'CHART': 'CHART',
            'TABLE': 'TABLE',
            'PICTURE': 'PICTURE',
            'MEDIA_CLIP': 'MEDIA',
            'ORG_CHART': 'ORG_CHART',
            'DATE': 'DATE',
            'FOOTER': 'FOOTER',
            'SLIDE_NUMBER': 'SLIDE_NUMBER',
            'HEADER': 'HEADER',
        }

        # Track placeholder counts to match naming
        placeholder_counts = {}

        # Fill placeholders
        try:
            for placeholder in slide.placeholders:
                if placeholder.has_text_frame:
                    try:
                        ph_idx = placeholder.placeholder_format.idx
                        ph_type = placeholder.placeholder_format.type

                        # Determine field name using same logic as extraction
                        field_name = f"PLACEHOLDER_{ph_idx}"
                        if ph_type:
                            type_str = str(ph_type)
                            base_name = type_str.split(' ')[0].split('(')[0].upper()
                            field_name = placeholder_name_map.get(base_name, base_name)

                        # Handle duplicate names
                        if field_name in placeholder_counts:
                            placeholder_counts[field_name] += 1
                            field_name = f"{field_name}_{placeholder_counts[field_name]}"
                        else:
                            placeholder_counts[field_name] = 1

                        if field_name in fields:
                            placeholder.text_frame.text = fields[field_name]
                    except Exception:
                        pass
        except Exception:
            pass

        # Fill text patterns
        try:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            new_text = run.text
                            for match in field_pattern.findall(run.text):
                                if match in fields:
                                    new_text = new_text.replace(f"{{{match}}}", fields[match])
                            if new_text != run.text:
                                run.text = new_text
        except Exception:
            pass

    def fill_presentation_fields(self, pptx_path: str, fields_by_slide: Dict[int, Dict[str, str]]) -> str:
        """
        Fill fields in a presentation and return path to the modified file.
        fields_by_slide: {slide_index: {field_name: value}}
        """
        if not os.path.exists(pptx_path):
            return pptx_path

        try:
            prs = Presentation(pptx_path)

            for slide_idx, fields in fields_by_slide.items():
                if 0 <= slide_idx < len(prs.slides):
                    self.fill_slide_fields(prs.slides[slide_idx], fields)

            # Save to temp file
            temp_path = tempfile.mktemp(suffix='.pptx')
            prs.save(temp_path)
            return temp_path

        except Exception as e:
            print(f"Error filling fields in {pptx_path}: {e}")
            return pptx_path

    def merge_liturgy_v2(self, liturgy: Liturgy) -> str:
        """
        Create a merged presentation from liturgy sections (v2 format).
        Returns path to the merged file.
        """
        if not liturgy.sections:
            # Fall back to v1 merge if no sections
            return self.merge_liturgy(liturgy)

        # Collect all source files and slide indices to copy
        slides_to_copy = []  # List of (pptx_path, slide_indices, fields_by_index)

        for section in liturgy.sections:
            if section.is_song:
                # For song sections, each slide entry represents a song
                # Group slides by source_path to avoid duplicating slides from same file
                processed_sources = set()
                for slide in section.slides:
                    if slide.is_stub:
                        stub_path = self._create_stub_presentation(slide.title or section.name)
                        slides_to_copy.append((stub_path, [0], {0: slide.fields}))
                    elif slide.source_path and os.path.exists(slide.source_path):
                        # Check if we already processed this source file
                        if slide.source_path in processed_sources:
                            continue
                        processed_sources.add(slide.source_path)

                        slide_count = self.get_slide_count(slide.source_path)
                        slides_to_copy.append((
                            slide.source_path,
                            list(range(slide_count)),
                            {i: slide.fields for i in range(slide_count)},
                        ))
            else:
                # For regular sections, process each slide individually
                for slide in section.slides:
                    if slide.is_stub:
                        stub_path = self._create_stub_presentation(slide.title or section.name)
                        slides_to_copy.append((stub_path, [0], {0: slide.fields}))
                    elif slide.source_path and os.path.exists(slide.source_path):
                        slides_to_copy.append((
                            slide.source_path,
                            [slide.slide_index],
                            {slide.slide_index: slide.fields},
                        ))

        if not slides_to_copy:
            # No slides to copy, create empty presentation
            logger.info("merge_liturgy_v2: No slides to copy, creating empty presentation")
            try:
                prs = Presentation()
                prs.slide_width = Inches(13.333)
                prs.slide_height = Inches(7.5)
                temp_file = tempfile.NamedTemporaryFile(suffix='.pptx', delete=False)
                temp_path = temp_file.name
                temp_file.close()
                prs.save(temp_path)
                logger.debug(f"Empty presentation saved to: {temp_path}")
                return temp_path
            except Exception as e:
                logger.error(f"Failed to create empty presentation: {e}", exc_info=True)
                raise

        # Try VBScript method (Windows with PowerPoint installed)
        if os.name == 'nt':
            try:
                # For VBScript, we need to fill fields after merging
                result_path = self._merge_with_vbscript(
                    [(path, indices) for path, indices, _ in slides_to_copy]
                )

                # Apply fields to merged result
                # Build a mapping of result slide index -> fields
                result_slide_idx = 0
                fields_for_result = {}
                for _, slide_indices, fields_by_index in slides_to_copy:
                    for src_idx in slide_indices:
                        if src_idx in fields_by_index and fields_by_index[src_idx]:
                            fields_for_result[result_slide_idx] = fields_by_index[src_idx]
                        result_slide_idx += 1

                if fields_for_result:
                    result_path = self.fill_presentation_fields(result_path, fields_for_result)

                return result_path
            except Exception as e:
                print(f"VBScript merge failed: {e}, falling back to python-pptx")

        # Fallback to python-pptx method with field filling
        return self._merge_with_pptx_v2(slides_to_copy)

    def _merge_with_pptx_v2(
        self,
        slides_to_copy: List[Tuple[str, List[int], Dict[int, Dict[str, str]]]]
    ) -> str:
        """
        Merge presentations using python-pptx with field filling support.
        slides_to_copy: List of (pptx_path, slide_indices, fields_by_index)
        """
        if not slides_to_copy:
            logger.info("_merge_with_pptx_v2: No slides to copy, creating empty presentation")
            try:
                prs = Presentation()
                prs.slide_width = Inches(13.333)
                prs.slide_height = Inches(7.5)
                temp_file = tempfile.NamedTemporaryFile(suffix='.pptx', delete=False)
                temp_path = temp_file.name
                temp_file.close()
                prs.save(temp_path)
                logger.debug(f"Empty presentation saved to: {temp_path}")
                return temp_path
            except Exception as e:
                logger.error(f"Failed to create empty presentation: {e}", exc_info=True)
                raise

        # Use first presentation as base
        first_path, first_indices, first_fields = slides_to_copy[0]
        logger.info(f"Merging {len(slides_to_copy)} presentations using python-pptx v2")
        try:
            prs = Presentation(first_path)
            logger.debug(f"Loaded base presentation: {first_path}")
        except Exception as e:
            logger.error(f"Failed to load base presentation {first_path}: {e}", exc_info=True)
            raise

        # Remove slides we don't want from the first presentation
        all_indices = set(range(len(prs.slides)))
        indices_to_remove = all_indices - set(first_indices)

        # Remove in reverse order to maintain indices
        for idx in sorted(indices_to_remove, reverse=True):
            rId = prs.slides._sldIdLst[idx].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[idx]

        # Fill fields in the remaining slides from first presentation
        for i, slide_idx in enumerate(first_indices):
            if slide_idx in first_fields and i < len(prs.slides):
                self.fill_slide_fields(prs.slides[i], first_fields[slide_idx])

        # Add slides from other presentations
        for source_path, slide_indices, fields_by_index in slides_to_copy[1:]:
            if not os.path.exists(source_path):
                logger.warning(f"Source presentation not found, skipping: {source_path}")
                continue

            try:
                source_prs = Presentation(source_path)
                logger.debug(f"Loading source presentation: {source_path}")
                for slide_idx in slide_indices:
                    new_slide = self._copy_slide_from_source(prs, source_prs, slide_idx)
                    if new_slide and slide_idx in fields_by_index:
                        self.fill_slide_fields(new_slide, fields_by_index[slide_idx])
            except Exception as e:
                logger.error(f"Failed to load/copy from {source_path}: {e}", exc_info=True)
                continue

        temp_file = tempfile.NamedTemporaryFile(suffix='.pptx', delete=False)
        temp_path = temp_file.name
        temp_file.close()
        prs.save(temp_path)
        return temp_path
