"""Bible slide generation service.

Creates PowerPoint slides containing Bible text in multiple translations,
laid out in a grid (one column per translation).  Each slide shows a subset
of the requested verse range (determined by a text-fitting algorithm) and
includes a QR code linking to the full passage on YouVersion.

Verse alignment between translations is done by row-index (0-based position
in the fetched verse list) rather than verse number, so per-translation
reference overrides that cover the same passage but with different numbering
will still line up correctly.
"""

import os
import re
import tempfile
from dataclasses import dataclass, field
from io import BytesIO
from typing import Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt, Emu

from ..logging_config import get_logger
from .bible_service import (
    BibleReference,
    BibleService,
    BibleTranslation,
    BibleVerse,
    parse_reference,
    parse_references,
)
from .google_translate_service import is_rtl

logger = get_logger("bible_slide_service")


# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------

@dataclass
class BibleSlideConfig:
    """Configuration for Bible slide generation."""

    font_name: str = "Calibri"
    font_size: int = 12           # points
    max_chars_per_slide: int = 500  # total characters before starting a new slide
    slide_width: float = 10.0     # inches (standard 4:3)
    slide_height: float = 7.5     # inches
    show_verse_numbers: bool = True
    title_font_size: int = 18     # points
    header_font_size: int = 11    # points (translation name row)
    title_bold: bool = True
    bg_color: Optional[Tuple[int, int, int]] = None   # None = white
    title_color: Tuple[int, int, int] = (0, 0, 0)     # black
    text_color: Tuple[int, int, int] = (0, 0, 0)      # black
    header_color: Tuple[int, int, int] = (80, 80, 80) # dark grey
    qr_size: float = 1.0          # inches (square)
    column_gap: float = 0.1       # inches between columns
    margin_left: float = 0.3      # inches
    margin_right: float = 0.3     # inches
    margin_top: float = 0.25      # inches
    margin_bottom: float = 0.25   # inches


# Vertical layout constants (as fractions of slide height, resolved at runtime)
_TITLE_HEIGHT_FRAC = 0.11   # 11 % of slide height for title bar
_HEADER_HEIGHT_FRAC = 0.07  # 7 % for translation-name header row

# Sentence boundary pattern: split after . ! ? followed by whitespace
_SENTENCE_SPLIT_RE = re.compile(r'(?<=[.!?])\s+')


# ---------------------------------------------------------------------------
# TranslationSlot – one column in the grid
# ---------------------------------------------------------------------------

@dataclass
class TranslationSlot:
    """A single translation column: translation metadata + fetched verse list.

    The *reference_override* (if set) is the reference actually used to fetch
    verses for this translation; it may differ from the main reference to
    accommodate verse-numbering differences between traditions.
    """
    translation: BibleTranslation
    verses: List[BibleVerse]           # in order fetched
    reference_override: Optional[str] = None  # if different from main ref


# ---------------------------------------------------------------------------
# BibleSlideService
# ---------------------------------------------------------------------------

class BibleSlideService:
    """Creates PowerPoint presentations from Bible text fetched via BibleService."""

    def __init__(self, bible_service: Optional[BibleService] = None) -> None:
        self._bible = bible_service or BibleService()

    # ------------------------------------------------------------------
    # Public
    # ------------------------------------------------------------------

    def create_slides(
        self,
        reference_str: str,
        translation_ids: List[int],
        config: Optional[BibleSlideConfig] = None,
        reference_overrides: Optional[Dict[int, str]] = None,
        template_path: Optional[str] = None,
    ) -> str:
        """Generate a temporary PPTX file with Bible text slides.

        Args:
            reference_str: Human reference string, e.g. "John 3:16-21".
            translation_ids: List of YouVersion version IDs (max 6).
            config: Slide configuration.  Defaults used if not provided.
            reference_overrides: Map from version_id → reference string to
                use instead of *reference_str* for that specific translation.
                Useful when different traditions use different verse numbering
                for the same passage.

        Returns:
            Absolute path to the generated temporary PPTX file.
            Caller is responsible for moving / deleting it.
        """
        if config is None:
            config = BibleSlideConfig()

        if not translation_ids:
            raise ValueError("At least one translation ID is required.")

        translation_ids = translation_ids[:6]
        reference_overrides = reference_overrides or {}

        # 1. Parse the main reference (supports comma-separated multiple refs)
        main_refs = parse_references(reference_str)
        is_multi_ref = len(main_refs) > 1
        display_label = (
            reference_str.strip() if is_multi_ref else main_refs[0].display_str
        )
        logger.info(
            "Creating Bible slides for: %s, translations: %s",
            display_label, translation_ids,
        )

        # 2. Build translation catalog lookup
        all_trans = self._bible.get_builtin_translations()
        id_to_trans: Dict[int, BibleTranslation] = {t.id: t for t in all_trans}

        # 3. Fetch verses per translation
        slots: List[TranslationSlot] = []
        for vid in translation_ids:
            trans = id_to_trans.get(vid, BibleTranslation(
                id=vid, abbreviation=str(vid), name=str(vid), language="", language_name=""
            ))
            override_str = reference_overrides.get(vid)
            try:
                if override_str:
                    override_refs = parse_references(override_str)
                    verses = self._bible.get_verses_multi(override_refs, vid)
                elif is_multi_ref:
                    verses = self._bible.get_verses_multi(main_refs, vid)
                else:
                    verses = self._bible.get_verses(main_refs[0], vid)
                if not verses:
                    logger.warning("No verses for version %s", vid)
                    verses = [BibleVerse(verse_num=0, text="(No text available)")]
            except Exception as exc:
                logger.error("Failed to fetch verses for version %s: %s", vid, exc, exc_info=True)
                verses = [BibleVerse(verse_num=0, text=f"(Error: {exc})")]

            slots.append(TranslationSlot(
                translation=trans,
                verses=verses,
                reference_override=override_str,
            ))

        # 4. QR URL (use first translation, first reference)
        qr_url = self._bible.get_youversion_url(main_refs[0], translation_ids[0])

        # Use a synthetic BibleReference for the slide title
        title_ref = main_refs[0]

        # 5. Build slides
        return self._build_presentation(
            title_ref, slots, qr_url, config, display_label,
            template_path=template_path,
        )

    def fetch_verses_for_slot(
        self,
        translation_id: int,
        reference_str: str,
    ) -> List[BibleVerse]:
        """Fetch verses for a single translation, used by the preview pane."""
        refs = parse_references(reference_str)
        if len(refs) == 1:
            return self._bible.get_verses(refs[0], translation_id)
        return self._bible.get_verses_multi(refs, translation_id)

    # ------------------------------------------------------------------
    # Private – presentation builder
    # ------------------------------------------------------------------

    def _build_presentation(
        self,
        main_ref: BibleReference,
        slots: List[TranslationSlot],
        qr_url: str,
        config: BibleSlideConfig,
        display_label: Optional[str] = None,
        template_path: Optional[str] = None,
    ) -> str:
        prs = Presentation()
        prs.slide_width = Inches(config.slide_width)
        prs.slide_height = Inches(config.slide_height)

        geom = _SlideGeometry(config)
        n_cols = len(slots)

        # Align by row-index: each "row" is the i-th verse in the fetched list.
        # This handles verse-numbering discrepancies between translations naturally.
        max_rows = max((len(s.verses) for s in slots), default=0)
        if max_rows == 0:
            raise ValueError("No verse data retrieved for any translation.")

        # verse_lookup[(col_idx, row_idx)] = text shown in that cell
        verse_lookup: Dict[Tuple[int, int], str] = {}
        for col_idx, slot in enumerate(slots):
            for row_idx, v in enumerate(slot.verses):
                text = v.text
                if config.show_verse_numbers and v.verse_num != 0:
                    text = f"{v.verse_num}\u00a0{text}"  # non-breaking space
                verse_lookup[(col_idx, row_idx)] = text

        all_row_indices = list(range(max_rows))

        # Expand any row whose text (in any column) exceeds max_chars_per_slide
        # by splitting at sentence boundaries.  This produces sub-rows that are
        # aligned across all columns (shorter translations get empty sub-rows).
        verse_lookup, all_row_indices = _expand_long_verses(
            verse_lookup, all_row_indices, n_cols, config.max_chars_per_slide
        )

        # Determine which rows go on which slide
        row_groups = self._group_rows_for_slides(
            all_row_indices=all_row_indices,
            verse_lookup=verse_lookup,
            n_translations=n_cols,
            geom=geom,
            config=config,
        )

        # Detect RTL columns
        rtl_flags = [is_rtl(s.translation.language) for s in slots]

        # Generate QR image once
        qr_image: Optional[BytesIO] = _generate_qr(qr_url)

        # Try template-based generation if a template is available
        if template_path and os.path.exists(template_path):
            try:
                return self._build_from_template(
                    template_path, main_ref, slots, row_groups,
                    verse_lookup, rtl_flags, qr_image, config,
                    display_label=display_label,
                )
            except Exception as exc:
                logger.warning(
                    "Template-based generation failed, falling back to default: %s",
                    exc, exc_info=True,
                )

        # --- From-scratch fallback ---
        blank_layout = prs.slide_layouts[6]
        slide_label = display_label if display_label else main_ref.display_str

        for group_idx, group_rows in enumerate(row_groups):
            slide = prs.slides.add_slide(blank_layout)

            if config.bg_color:
                _fill_slide_bg(slide, config.bg_color)

            # Title
            title_text = slide_label
            if len(row_groups) > 1:
                # Show verse sub-range for this slide using first column
                vr_start = slots[0].verses[group_rows[0]].verse_num if group_rows and slots[0].verses else "?"
                vr_end = slots[0].verses[group_rows[-1]].verse_num if group_rows and slots[0].verses else "?"
                if vr_start != vr_end:
                    title_text += f"  ({main_ref.chapter}:{vr_start}\u2013{vr_end})"
            _add_text_box(
                slide,
                x=Inches(config.margin_left),
                y=Inches(config.margin_top),
                w=Inches(config.slide_width - config.margin_left - config.margin_right),
                h=Inches(geom.title_height_in),
                text=title_text,
                font_name=config.font_name,
                font_size=config.title_font_size,
                bold=config.title_bold,
                color=config.title_color,
                align=PP_ALIGN.LEFT,
            )

            # Translation headers
            for col_idx, slot in enumerate(slots):
                col_x = Inches(geom.column_x(col_idx, n_cols))
                col_w = Inches(geom.column_width(n_cols))
                header_y = Inches(
                    config.margin_top + geom.title_height_in + geom.header_gap_in
                )
                ref_label = f" [{slot.reference_override}]" if slot.reference_override else ""
                header_text = f"{slot.translation.abbreviation} – {slot.translation.name}{ref_label}"
                _add_text_box(
                    slide,
                    x=col_x,
                    y=header_y,
                    w=col_w,
                    h=Inches(geom.header_height_in),
                    text=header_text,
                    font_name=config.font_name,
                    font_size=config.header_font_size,
                    bold=False,
                    color=config.header_color,
                    align=PP_ALIGN.LEFT,
                )

            # Verse text columns
            text_top = (
                config.margin_top
                + geom.title_height_in
                + geom.header_gap_in
                + geom.header_height_in
                + geom.text_top_gap_in
            )

            for col_idx, slot in enumerate(slots):
                col_x = Inches(geom.column_x(col_idx, n_cols))
                col_w = Inches(geom.column_width(n_cols))
                text_h = Inches(geom.text_area_height_in(config))

                lines = []
                for row_idx in group_rows:
                    lines.append(verse_lookup.get((col_idx, row_idx), ""))
                full_text = "\n".join(lines)

                _add_text_box(
                    slide,
                    x=col_x,
                    y=Inches(text_top),
                    w=col_w,
                    h=text_h,
                    text=full_text,
                    font_name=config.font_name,
                    font_size=config.font_size,
                    bold=False,
                    color=config.text_color,
                    align=PP_ALIGN.RIGHT if rtl_flags[col_idx] else PP_ALIGN.LEFT,
                    word_wrap=True,
                    rtl=rtl_flags[col_idx],
                )

            # QR code
            if qr_image:
                qr_x = Inches(config.slide_width - config.margin_right - config.qr_size)
                qr_y = Inches(config.slide_height - config.margin_bottom - config.qr_size)
                qr_image.seek(0)
                slide.shapes.add_picture(
                    qr_image,
                    left=qr_x,
                    top=qr_y,
                    width=Inches(config.qr_size),
                    height=Inches(config.qr_size),
                )

        fd, tmp_path = tempfile.mkstemp(suffix=".pptx", prefix="bible_slides_")
        os.close(fd)
        prs.save(tmp_path)
        logger.info("Bible slides saved to: %s (%d slides)", tmp_path, len(row_groups))
        return tmp_path

    # ------------------------------------------------------------------
    # Private – text fitting
    # ------------------------------------------------------------------

    def _group_rows_for_slides(
        self,
        all_row_indices: List[int],
        verse_lookup: Dict[Tuple[int, int], str],
        n_translations: int,
        geom: "_SlideGeometry",
        config: BibleSlideConfig,
    ) -> List[List[int]]:
        """Bin rows into slides using a character-count limit.

        Accumulates rows until adding the next row would exceed
        max_chars_per_slide, then starts a new slide.  The limit is compared
        against the *longest* translation column for each row so the widest
        text drives the split decision.
        """
        max_chars = config.max_chars_per_slide
        groups: List[List[int]] = []
        current_group: List[int] = []
        current_chars: int = 0

        for row_idx in all_row_indices:
            row_chars = max(
                len(verse_lookup.get((t, row_idx), ""))
                for t in range(n_translations)
            )
            if current_group and current_chars + row_chars > max_chars:
                groups.append(current_group)
                current_group = [row_idx]
                current_chars = row_chars
            else:
                current_group.append(row_idx)
                current_chars += row_chars

        if current_group:
            groups.append(current_group)
        return groups or [[]]

    # ------------------------------------------------------------------
    # Private – template-based builder
    # ------------------------------------------------------------------

    def _build_from_template(
        self,
        template_path: str,
        main_ref: BibleReference,
        slots: List[TranslationSlot],
        row_groups: List[List[int]],
        verse_lookup: Dict[Tuple[int, int], str],
        rtl_flags: List[bool],
        qr_image: Optional[BytesIO],
        config: BibleSlideConfig,
        display_label: Optional[str] = None,
    ) -> str:
        prs = Presentation(template_path)
        n_cols = len(slots)
        slide_label = display_label if display_label else main_ref.display_str

        # Build map: column_count → slide index in the template
        template_map = _build_template_map(prs)
        if not template_map:
            raise ValueError("No {CONTENT_N} placeholders found in template")

        logger.info(
            "Bible template map: %s (need %d columns)",
            template_map, n_cols,
        )

        # Determine how to split columns across template slides
        split_plan = _plan_column_splits(n_cols, sorted(template_map.keys(), reverse=True))

        n_template_slides = len(prs.slides)

        for group_idx, group_rows in enumerate(row_groups):
            # Build title text
            title_text = slide_label
            if len(row_groups) > 1:
                vr_start = slots[0].verses[group_rows[0]].verse_num if group_rows and slots[0].verses else "?"
                vr_end = slots[0].verses[group_rows[-1]].verse_num if group_rows and slots[0].verses else "?"
                if vr_start != vr_end:
                    title_text += f"  ({main_ref.chapter}:{vr_start}\u2013{vr_end})"

            for split_idx, (tpl_size, col_indices) in enumerate(split_plan):
                tpl_slide_idx, _chars_hint = template_map[tpl_size]
                template_slide = prs.slides[tpl_slide_idx]
                new_slide = _clone_slide(prs, template_slide)

                # Build fields dict: map template {CONTENT_1} etc. to actual column data
                fields: Dict[str, str] = {"TITLE": title_text}
                for local_idx, global_col in enumerate(col_indices):
                    slot = slots[global_col]
                    n = local_idx + 1
                    ref_label = f" [{slot.reference_override}]" if slot.reference_override else ""
                    fields[f"HEADER_{n}"] = f"{slot.translation.abbreviation} \u2013 {slot.translation.name}{ref_label}"

                    lines = []
                    for row_idx in group_rows:
                        lines.append(verse_lookup.get((global_col, row_idx), ""))
                    fields[f"CONTENT_{n}"] = "\n".join(lines)

                _fill_template_fields(new_slide, fields)

                # QR: only on the last slide of each verse group
                is_last_split = split_idx == len(split_plan) - 1
                if is_last_split and qr_image:
                    _replace_qr_placeholder(new_slide, qr_image)

                # Remove shapes still containing unfilled {FIELD} patterns
                _remove_unfilled_placeholders(new_slide)

        # Remove original template slides (iterate in reverse to keep indices stable)
        xml_sldIdLst = prs.slides._sldIdLst
        slide_ids = list(xml_sldIdLst)
        for sid in slide_ids[:n_template_slides]:
            xml_sldIdLst.remove(sid)

        fd, tmp_path = tempfile.mkstemp(suffix=".pptx", prefix="bible_slides_")
        os.close(fd)
        prs.save(tmp_path)
        logger.info(
            "Bible slides (template) saved to: %s (%d slides)",
            tmp_path, len(prs.slides),
        )
        return tmp_path


# ---------------------------------------------------------------------------
# Template helpers
# ---------------------------------------------------------------------------

_FIELD_PATTERN = re.compile(r'\{([A-Za-z_][A-Za-z0-9_]*)\}')
_CONTENT_N_PATTERN = re.compile(r'CONTENT_(\d+)', re.IGNORECASE)
_CHARS_HINT_PATTERN = re.compile(r'\{CHARS:(\d+)\}', re.IGNORECASE)


def _build_template_map(prs) -> Dict[int, Tuple[int, Optional[int]]]:
    """Scan all slides and return {column_count: (slide_index, suggested_chars)}.

    Column count is determined by the highest N in {CONTENT_N} on each slide.
    suggested_chars comes from a {CHARS:NNN} placeholder on the same slide, or None.
    """
    result: Dict[int, Tuple[int, Optional[int]]] = {}
    for slide_idx, slide in enumerate(prs.slides):
        max_n = 0
        suggested: Optional[int] = None
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                for field_name in _FIELD_PATTERN.findall(text):
                    m = _CONTENT_N_PATTERN.match(field_name)
                    if m:
                        max_n = max(max_n, int(m.group(1)))
                m_chars = _CHARS_HINT_PATTERN.search(text)
                if m_chars:
                    suggested = int(m_chars.group(1))
        if max_n > 0:
            # First slide wins for each column count
            if max_n not in result:
                result[max_n] = (slide_idx, suggested)
    return result


def _plan_column_splits(
    n_cols: int,
    available_sizes: List[int],
) -> List[Tuple[int, List[int]]]:
    """Greedily split n_cols across available template sizes.

    Returns list of (template_size, [column_indices]) tuples.
    E.g. n_cols=4, available=[3,2,1] → [(3, [0,1,2]), (1, [3])]
    """
    if not available_sizes:
        raise ValueError("No template slides available")

    plan: List[Tuple[int, List[int]]] = []
    remaining = list(range(n_cols))

    while remaining:
        # Find the largest template that fits without exceeding remaining columns
        best = None
        for size in available_sizes:
            if size <= len(remaining):
                best = size
                break
        if best is None:
            # All templates are larger than remaining; use smallest available
            best = available_sizes[-1]

        chunk = remaining[:best]
        remaining = remaining[best:]
        plan.append((best, chunk))

    return plan


def get_template_chars_hint(template_path: str, n_cols: int) -> Optional[int]:
    """Return minimum suggested chars-per-slide for *n_cols* translations.

    Reads *template_path*, builds the template map, determines which template
    slides will be used via _plan_column_splits, and returns the minimum
    {CHARS:NNN} value across those slides.  Returns None if no slide carries a
    hint or the template cannot be read.
    """
    try:
        prs = Presentation(template_path)
    except Exception as exc:
        logger.warning("get_template_chars_hint: could not open %s: %s", template_path, exc)
        return None
    template_map = _build_template_map(prs)
    if not template_map:
        return None
    available_sizes = sorted(template_map.keys(), reverse=True)
    try:
        split_plan = _plan_column_splits(n_cols, available_sizes)
    except ValueError:
        return None
    hints = [template_map[size][1] for size, _ in split_plan if template_map[size][1] is not None]
    return min(hints) if hints else None


def _clone_slide(prs, template_slide):
    """Clone a slide by deep-copying its XML into a new slide."""
    from copy import deepcopy
    from lxml import etree

    layout = template_slide.slide_layout
    new_slide = prs.slides.add_slide(layout)

    # Clear any auto-generated shapes from the layout
    spTree = new_slide.shapes._spTree
    for child in list(spTree):
        tag = etree.QName(child.tag).localname
        if tag in ('sp', 'pic', 'grpSp', 'graphicFrame', 'cxnSp'):
            spTree.remove(child)

    # Deep copy all shapes from template
    for child in template_slide.shapes._spTree:
        tag = etree.QName(child.tag).localname
        if tag in ('sp', 'pic', 'grpSp', 'graphicFrame', 'cxnSp'):
            spTree.append(deepcopy(child))

    # Copy background
    template_bg = template_slide.background._element
    new_bg = new_slide.background._element
    for child in list(new_bg):
        new_bg.remove(child)
    for child in template_bg:
        new_bg.append(deepcopy(child))

    return new_slide


def _fill_template_fields(slide, fields: Dict[str, str]) -> None:
    """Replace {FIELD_NAME} patterns in all text runs, preserving formatting."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    new_text = run.text
                    for match in _FIELD_PATTERN.findall(run.text):
                        if match in fields:
                            new_text = new_text.replace(f"{{{match}}}", fields[match])
                    if new_text != run.text:
                        run.text = new_text


def _replace_qr_placeholder(slide, qr_image: BytesIO) -> bool:
    """Find shape containing {QR}, replace with QR image at same position/size."""
    for shape in list(slide.shapes):
        if shape.has_text_frame and "{QR}" in shape.text_frame.text:
            left, top = shape.left, shape.top
            width, height = shape.width, shape.height
            shape._element.getparent().remove(shape._element)
            qr_image.seek(0)
            slide.shapes.add_picture(
                qr_image, left=left, top=top,
                width=width, height=height,
            )
            return True
    return False


def _remove_unfilled_placeholders(slide) -> None:
    """Remove shapes that still contain unfilled {FIELD} or {CHARS:NNN} patterns."""
    for shape in list(slide.shapes):
        if shape.has_text_frame:
            text = shape.text_frame.text
            if _FIELD_PATTERN.search(text) or _CHARS_HINT_PATTERN.search(text):
                shape._element.getparent().remove(shape._element)


# ---------------------------------------------------------------------------
# Slide geometry helper
# ---------------------------------------------------------------------------

class _SlideGeometry:
    def __init__(self, config: BibleSlideConfig) -> None:
        self.c = config
        self.title_height_in: float = config.slide_height * _TITLE_HEIGHT_FRAC
        self.header_gap_in: float = 0.05
        self.header_height_in: float = config.slide_height * _HEADER_HEIGHT_FRAC
        self.text_top_gap_in: float = 0.05

    def column_width(self, n_cols: int) -> float:
        c = self.c
        total_gap = c.column_gap * (n_cols - 1) if n_cols > 1 else 0
        usable_w = c.slide_width - c.margin_left - c.margin_right - total_gap
        return max(0.5, usable_w / n_cols)

    def column_x(self, col_idx: int, n_cols: int) -> float:
        c = self.c
        return c.margin_left + col_idx * (self.column_width(n_cols) + c.column_gap)

    def text_area_height_in(self, config: BibleSlideConfig) -> float:
        used = (
            config.margin_top
            + self.title_height_in
            + self.header_gap_in
            + self.header_height_in
            + self.text_top_gap_in
            + config.margin_bottom
            + config.qr_size
        )
        return max(0.5, config.slide_height - used)


# ---------------------------------------------------------------------------
# Verse-splitting helpers
# ---------------------------------------------------------------------------

def _split_at_sentence_boundaries(text: str, max_chars: int) -> List[str]:
    """Split *text* into parts of at most *max_chars* characters each.

    Cuts only at sentence endings (.  !  ?) followed by whitespace.
    If no sentence boundary exists the whole text is returned as a single
    part (no mid-word cuts ever occur).
    """
    if len(text) <= max_chars:
        return [text]
    sentences = _SENTENCE_SPLIT_RE.split(text)
    parts: List[str] = []
    current = ""
    for s in sentences:
        candidate = (current + " " + s).strip() if current else s
        if len(candidate) <= max_chars:
            current = candidate
        else:
            if current:
                parts.append(current)
            current = s
    if current:
        parts.append(current)
    return parts or [text]


def _expand_long_verses(
    verse_lookup: Dict[Tuple[int, int], str],
    all_row_indices: List[int],
    n_cols: int,
    max_chars: int,
) -> Tuple[Dict[Tuple[int, int], str], List[int]]:
    """Expand verse rows that exceed *max_chars* into sentence-split sub-rows.

    When a verse in *any* translation column is longer than *max_chars*, the
    row is replaced by multiple sub-rows — one per sentence group.  Shorter
    translations in the same row are placed entirely in the first sub-row;
    remaining sub-rows are left empty so layout stays aligned across columns.
    """
    expanded: Dict[Tuple[int, int], str] = {}
    new_indices: List[int] = []
    new_row = 0

    for orig_row in all_row_indices:
        splits_per_col: Dict[int, List[str]] = {}
        max_parts = 1
        for col_idx in range(n_cols):
            text = verse_lookup.get((col_idx, orig_row), "")
            parts = _split_at_sentence_boundaries(text, max_chars)
            splits_per_col[col_idx] = parts
            max_parts = max(max_parts, len(parts))

        for sub in range(max_parts):
            for col_idx in range(n_cols):
                parts = splits_per_col[col_idx]
                expanded[(col_idx, new_row)] = parts[sub] if sub < len(parts) else ""
            new_indices.append(new_row)
            new_row += 1

    return expanded, new_indices


# ---------------------------------------------------------------------------
# python-pptx helpers
# ---------------------------------------------------------------------------

def _add_text_box(
    slide,
    x: Emu,
    y: Emu,
    w: Emu,
    h: Emu,
    text: str,
    font_name: str,
    font_size: int,
    bold: bool,
    color: Tuple[int, int, int],
    align: PP_ALIGN = PP_ALIGN.LEFT,
    word_wrap: bool = True,
    rtl: bool = False,
) -> None:
    """Add a text box with optional RTL paragraph direction."""
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap

    paragraphs = text.split("\n")
    for p_idx, para_text in enumerate(paragraphs):
        if p_idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align

        # Set RTL paragraph direction via XML
        if rtl:
            _set_paragraph_rtl(p)

        run = p.add_run()
        run.text = para_text
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = RGBColor(*color)


def _set_paragraph_rtl(paragraph) -> None:
    """Set paragraph direction to RTL via XML manipulation."""
    try:
        pPr = paragraph._p.get_or_add_pPr()
        pPr.set(qn("a:rtl"), "1")
    except Exception as exc:
        logger.debug("Could not set RTL on paragraph: %s", exc)


def _fill_slide_bg(slide, color: Tuple[int, int, int]) -> None:
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*color)


# ---------------------------------------------------------------------------
# QR code generation
# ---------------------------------------------------------------------------

def _generate_qr(url: str) -> Optional[BytesIO]:
    try:
        import qrcode  # type: ignore
        qr = qrcode.QRCode(
            version=None,
            error_correction=qrcode.constants.ERROR_CORRECT_M,
            box_size=10,
            border=2,
        )
        qr.add_data(url)
        qr.make(fit=True)
        img = qr.make_image(fill_color="black", back_color="white")
        buf = BytesIO()
        img.save(buf, format="PNG")
        buf.seek(0)
        return buf
    except ImportError:
        logger.warning("qrcode library not available; QR codes will not be generated.")
        return None
    except Exception as exc:
        logger.error("Failed to generate QR code: %s", exc, exc_info=True)
        return None
